# -*- coding: utf-8 -*-
"""PPT for AI solution achievements: industry solutions + scenario examples.
Visual style aligned with AI应用开发典型成果与场景汇报-v3.pptx.
"""
from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib.patches import Circle, FancyBboxPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

C = {
    "bg": "F4F6F8",
    "white": "FFFFFF",
    "ink": "1A2332",
    "muted": "5B6B7C",
    "line": "D5DCE4",
    "navy": "1B3A4B",
    "teal": "0F766E",
    "amber": "B45309",
}

OUT_DIR = Path(__file__).resolve().parent
ASSETS = OUT_DIR / "_ppt_assets"
ASSETS.mkdir(exist_ok=True)
OUT_PPTX = OUT_DIR / "AI应用解决方案成果积累和应用场景示例_bak.pptx"

plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False


def hex_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_run(run, text, size=14, bold=False, color=C["ink"], font="Microsoft YaHei"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = hex_rgb(color)
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = parse_xml(
        f'<a:ea xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{font}"/>'
    )
    rPr.append(ea)


def add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=C["ink"], align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size=size, bold=bold, color=color)
    return box


def add_para(tf, text, size=13, bold=False, color=C["ink"], space_before=6):
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    set_run(p.add_run(), text, size=size, bold=bold, color=color)
    return p


def fill_shape(shape, color_hex: str):
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(color_hex)
    shape.line.fill.background()


def rounded_rect(slide, left, top, width, height, fill=C["white"], line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_shape(shape, fill)
    if line:
        shape.line.color.rgb = hex_rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def set_slide_bg(slide, color_hex=C["bg"]):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    fill_shape(bg, color_hex)
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.4), color=C["teal"]):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_shape(s, color)
    return s


def page_header(slide, title: str, subtitle: str | None = None, section: str | None = None):
    set_slide_bg(slide)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
    fill_shape(strip, C["navy"])
    if section:
        add_textbox(slide, Inches(0.7), Inches(0.28), Inches(11.5), Inches(0.3), section, size=11, color=C["teal"], bold=True)
    accent_bar(slide, Inches(0.7), Inches(0.58 if section else 0.45), height=Inches(0.36))
    add_textbox(
        slide,
        Inches(0.95),
        Inches(0.5 if section else 0.38),
        Inches(11.5),
        Inches(0.45),
        title,
        size=24,
        bold=True,
        color=C["navy"],
    )
    if subtitle:
        add_textbox(slide, Inches(0.95), Inches(1.0), Inches(11.5), Inches(0.35), subtitle, size=13, color=C["muted"])


def footer(slide, page: int, total: int, light: bool = True):
    ink = C["muted"] if light else "94A3B8"
    add_textbox(slide, Inches(0.7), Inches(7.1), Inches(9), Inches(0.3), "AI 应用解决方案  ·  成果积累与应用场景示例", size=10, color=ink)
    add_textbox(slide, Inches(11.2), Inches(7.1), Inches(1.5), Inches(0.3), f"{page} / {total}", size=10, color=ink, align=PP_ALIGN.RIGHT)


def add_labeled_block(slide, left, top, width, height, label, body, label_color=C["amber"], accent=C["teal"]):
    rounded_rect(slide, left, top, width, height, fill=C["white"], line=C["line"])
    accent_bar(slide, left, top, width=Inches(0.1), height=height, color=accent)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.1), width - Inches(0.4), Inches(0.32), label, size=12, bold=True, color=label_color)
    box = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.42), width - Inches(0.4), height - Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in body:
        if first:
            set_run(tf.paragraphs[0].add_run(), line, size=12, color=C["ink"])
            first = False
        else:
            add_para(tf, line, size=12, color=C["ink"], space_before=4)
    return box


def add_detail_slide(prs, blank, page, total, title, pain_points, solutions, delivers, mapping, chart_path, accent=C["teal"]):
    s = prs.slides.add_slide(blank)
    page_header(
        s,
        title,
        "行业场景痛点 → 解决方案 → 可交付能力",
        section="PART 01  ·  行业方案详解",
    )
    w = Inches(6.35)
    add_labeled_block(s, Inches(0.5), Inches(1.35), w, Inches(1.7), "行业场景痛点", pain_points, label_color=C["amber"], accent=C["amber"])
    add_labeled_block(s, Inches(0.5), Inches(3.2), w, Inches(1.55), "解决方案", solutions, label_color=C["navy"], accent=accent)
    add_labeled_block(s, Inches(0.5), Inches(4.9), w, Inches(1.35), "可交付能力", delivers, label_color=C["teal"], accent=C["teal"])
    add_textbox(s, Inches(0.55), Inches(6.4), w, Inches(0.35), "本项目映射：" + mapping, size=12, color=C["muted"])
    rounded_rect(s, Inches(7.05), Inches(1.35), Inches(5.7), Inches(5.35), fill=C["white"], line=C["line"])
    s.shapes.add_picture(str(chart_path), Inches(7.2), Inches(1.5), height=Inches(5.05))
    footer(s, page, total)
    return s


def chart_overview(path: Path):
    fig, ax = plt.subplots(figsize=(12.4, 5.0), dpi=160)
    ax.set_xlim(0, 12.4)
    ax.set_ylim(0, 5.0)
    ax.axis("off")
    fig.patch.set_facecolor("#F4F6F8")
    ax.set_facecolor("#F4F6F8")
    layers = [
        (2.85, "#1B3A4B", "行业场景解决方案", "问答查数 · 综合分析 · 现场研判 · 文档整理 · 风险研判"),
        (1.15, "#0F766E", "应用场景示例（项目落地）", "AI问答 · 超温/专项分析 · 看图诊断 · 检修报告提取 · 风险行为研判"),
    ]
    for y, color, title, desc in layers:
        ax.add_patch(FancyBboxPatch((0.7, y), 11.0, 1.35, boxstyle="round,pad=0.02,rounding_size=0.12", facecolor=color, linewidth=0))
        ax.text(1.0, y + 0.85, title, color="white", fontsize=15, fontweight="bold")
        ax.text(1.0, y + 0.35, desc, color="#E2E8F0", fontsize=12)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def chart_industry(path: Path):
    fig, ax = plt.subplots(figsize=(12.4, 5.0), dpi=160)
    ax.set_xlim(0, 12.4)
    ax.set_ylim(0, 5.0)
    ax.axis("off")
    fig.patch.set_facecolor("#F4F6F8")
    ax.set_facecolor("#F4F6F8")
    # row1: 3 cards, row2: 2 cards centered
    row1 = [
        ("知识问答与数据查询一体化", "问知识 · 查业务数据"),
        ("多源业务综合分析", "多源取数 · 结构化报告"),
        ("多模态现场研判", "图像+历史+规范"),
    ]
    row2 = [
        ("业务文档结构化整理", "文档抽取 · 标准入库"),
        ("作业风险行为检测研判", "小模型检出 · 大模型二次研判"),
    ]
    colors = ["#0F766E", "#1B3A4B", "#0F766E", "#1B3A4B", "#0F766E"]
    w, gap = 3.5, 0.35
    x0 = 0.7
    for i, (t, d) in enumerate(row1):
        x = x0 + i * (w + gap)
        color = colors[i]
        ax.add_patch(FancyBboxPatch((x, 2.85), w, 1.55, boxstyle="round,pad=0.02,rounding_size=0.1", facecolor=color, linewidth=0))
        ax.text(x + w / 2, 3.85, t, ha="center", va="center", color="white", fontsize=13, fontweight="bold")
        ax.text(x + w / 2, 3.25, d, ha="center", va="center", color="#CCFBF1", fontsize=11)
    w2, gap2 = 4.6, 0.4
    x02 = 1.4
    for i, (t, d) in enumerate(row2):
        x = x02 + i * (w2 + gap2)
        color = colors[i + 3]
        ax.add_patch(FancyBboxPatch((x, 0.7), w2, 1.7, boxstyle="round,pad=0.02,rounding_size=0.1", facecolor=color, linewidth=0))
        ax.text(x + w2 / 2, 1.75, t, ha="center", va="center", color="white", fontsize=13, fontweight="bold")
        ax.text(x + w2 / 2, 1.15, d, ha="center", va="center", color="#E2E8F0", fontsize=11)
    ax.text(6.2, 4.7, "五大通用行业场景解决方案", ha="center", color="#1B3A4B", fontsize=15, fontweight="bold")
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def chart_solution_map(path: Path):
    """场景痛点 → 方案 → 交付 对照图."""
    fig, ax = plt.subplots(figsize=(12.4, 5.2), dpi=160)
    ax.set_xlim(0, 12.4)
    ax.set_ylim(0, 5.2)
    ax.axis("off")
    fig.patch.set_facecolor("#F4F6F8")
    ax.set_facecolor("#F4F6F8")

    headers = [(0.5, "行业场景痛点", "#B45309"), (4.4, "对应解决方案", "#0F766E"), (8.3, "可交付能力", "#1B3A4B")]
    for x, title, color in headers:
        ax.add_patch(FancyBboxPatch((x, 4.35), 3.6, 0.55, boxstyle="round,pad=0.02,rounding_size=0.08", facecolor=color, linewidth=0))
        ax.text(x + 1.8, 4.62, title, ha="center", va="center", color="white", fontsize=12, fontweight="bold")

    rows = [
        ("资料难查、数据难问", "知识问答与数据查询一体化", "统一问答 · 知识库 · 只读查数"),
        ("分析拼数慢、报告乱", "多源业务综合分析", "专题分析 · 结构化报告"),
        ("现场取证后信息割裂", "多模态现场研判", "图像+历史+规范合成"),
        ("报告录入慢、难入库", "业务文档结构化整理", "文档抽取 · 标准记录"),
        ("安监误报高、语义弱", "作业风险行为检测与研判", "小模型检出 · 大模型二次研判"),
    ]
    for i, (a, b, c) in enumerate(rows):
        y = 3.55 - i * 0.7
        for x, text, bg in ((0.5, a, "#FFF7ED"), (4.4, b, "#ECFDF5"), (8.3, c, "#F1F5F9")):
            ax.add_patch(
                FancyBboxPatch((x, y), 3.6, 0.58, boxstyle="round,pad=0.02,rounding_size=0.06", facecolor=bg, edgecolor="#D5DCE4", linewidth=1)
            )
            ax.text(x + 1.8, y + 0.29, text, ha="center", va="center", color="#1A2332", fontsize=10)
        ax.annotate("", xy=(4.35, y + 0.29), xytext=(4.15, y + 0.29), arrowprops=dict(arrowstyle="-|>", color="#94A3B8", lw=1.2))
        ax.annotate("", xy=(8.25, y + 0.29), xytext=(8.05, y + 0.29), arrowprops=dict(arrowstyle="-|>", color="#94A3B8", lw=1.2))

    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def chart_detail_pipeline(path: Path, title: str, steps: list[str], accent: str = "#0F766E"):
    """Single-scenario solve path for detail slides (vertical-friendly horizontal flow)."""
    fig, ax = plt.subplots(figsize=(6.4, 5.4), dpi=160)
    n = len(steps)
    ax.set_xlim(0, 6.4)
    ax.set_ylim(0, 5.4)
    ax.axis("off")
    fig.patch.set_facecolor("#F4F6F8")
    ax.set_facecolor("#F4F6F8")

    ax.add_patch(FancyBboxPatch((0.25, 4.7), 5.9, 0.5, boxstyle="round,pad=0.02,rounding_size=0.08", facecolor=accent, linewidth=0))
    ax.text(3.2, 4.95, title, ha="center", va="center", color="white", fontsize=13, fontweight="bold")

    box_h = 0.72
    gap = 0.18
    top_y = 3.85
    for i, step in enumerate(steps):
        y = top_y - i * (box_h + gap)
        ax.add_patch(
            FancyBboxPatch((0.7, y), 5.0, box_h, boxstyle="round,pad=0.02,rounding_size=0.1", facecolor="white", edgecolor=accent, linewidth=1.8)
        )
        ax.add_patch(Circle((1.25, y + box_h / 2), 0.22, color=accent))
        ax.text(1.25, y + box_h / 2, f"{i+1}", ha="center", va="center", color="white", fontsize=11, fontweight="bold")
        ax.text(1.7, y + box_h / 2, step, ha="left", va="center", color="#1A2332", fontsize=12)
        if i < n - 1:
            ax.annotate(
                "",
                xy=(3.2, y - gap),
                xytext=(3.2, y),
                arrowprops=dict(arrowstyle="-|>", color="#94A3B8", lw=1.4),
            )

    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def chart_solution_flows(path: Path):
    """三条代表性解决路径示意."""
    fig, ax = plt.subplots(figsize=(12.4, 5.0), dpi=160)
    ax.set_xlim(0, 12.4)
    ax.set_ylim(0, 5.0)
    ax.axis("off")
    fig.patch.set_facecolor("#F4F6F8")
    ax.set_facecolor("#F4F6F8")

    flows = [
        (3.55, "#0F766E", "问答查数", ["用户提问", "意图分流", "知识检索 / 查数", "流式回答"]),
        (2.05, "#1B3A4B", "综合分析", ["分析需求", "取数计划", "并行查数+知识", "结构化报告"]),
        (0.55, "#0F766E", "风险研判", ["小模型检出", "上送图像/目标", "大模型二次研判", "确认并分级"]),
    ]
    for y, color, title, steps in flows:
        ax.add_patch(FancyBboxPatch((0.4, y), 1.6, 1.1, boxstyle="round,pad=0.02,rounding_size=0.08", facecolor=color, linewidth=0))
        ax.text(1.2, y + 0.55, title, ha="center", va="center", color="white", fontsize=12, fontweight="bold")
        for i, step in enumerate(steps):
            x = 2.3 + i * 2.45
            ax.add_patch(
                FancyBboxPatch((x, y + 0.15), 2.2, 0.8, boxstyle="round,pad=0.02,rounding_size=0.08", facecolor="white", edgecolor=color, linewidth=1.5)
            )
            ax.text(x + 1.1, y + 0.55, step, ha="center", va="center", color="#1A2332", fontsize=11)
            if i < len(steps) - 1:
                ax.annotate("", xy=(x + 2.25, y + 0.55), xytext=(x + 2.15, y + 0.55), arrowprops=dict(arrowstyle="-|>", color="#94A3B8", lw=1.3))
    ax.text(6.2, 4.75, "代表性解决路径示意（详解节选）", ha="center", color="#1B3A4B", fontsize=14, fontweight="bold")
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def chart_scenes(path: Path):
    fig, ax = plt.subplots(figsize=(12.2, 4.8), dpi=160)
    ax.set_xlim(0, 12.2)
    ax.set_ylim(0, 4.8)
    ax.axis("off")
    fig.patch.set_facecolor("#F4F6F8")
    ax.set_facecolor("#F4F6F8")
    groups = [
        (0.4, 2.5, 3.7, 1.9, "#0F766E", "日常问询", ["AI问答", "独立自然语言查数", "知识库运维"]),
        (4.25, 2.5, 3.7, 1.9, "#1B3A4B", "分析研判", ["超温 / 检修策略 / 四管", "泄爆溯源分析", "分析报告结构化输出"]),
        (8.1, 2.5, 3.7, 1.9, "#0F766E", "现场与文档 / 安监", ["看图诊断（缺陷/泄爆）", "检修报告结构化提取", "风险行为检测与研判"]),
        (2.0, 0.3, 8.2, 1.8, "#334155", "效果共性", ["提效：开口驱动系统", "提质：有依据 · 结构统一", "可控：可确认 · 只读查数 · 可复制"]),
    ]
    for x, y, w, h, color, title, lines in groups:
        ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.02,rounding_size=0.12", facecolor=color, linewidth=0))
        ax.text(x + 0.25, y + h - 0.35, title, color="white", fontsize=13, fontweight="bold")
        for i, line in enumerate(lines):
            ax.text(x + 0.25, y + h - 0.75 - i * 0.38, "·  " + line, color="#E2E8F0", fontsize=11)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


PHOTO_DEFECT = ASSETS / "photo_defect.png"
PHOTO_RISK = ASSETS / "photo_risk_site.png"


def _load_rgb(path: Path):
    from PIL import Image
    import numpy as np

    return np.asarray(Image.open(path).convert("RGB"))


def _browser_frame(ax, x=0.08, y=0.08, w=7.64, h=5.64, url="https://ai-app.platform/apps/demo"):
    """More realistic browser chrome; returns content box (cx, cy, cw, ch)."""
    # outer shadow
    ax.add_patch(FancyBboxPatch((x + 0.04, y - 0.03), w, h, boxstyle="round,pad=0.01,rounding_size=0.1", facecolor="#CBD5E1", linewidth=0, zorder=0))
    # window
    ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.01,rounding_size=0.1", facecolor="#FFFFFF", edgecolor="#94A3B8", linewidth=1.2, zorder=1))
    # title bar
    ax.add_patch(FancyBboxPatch((x, y + h - 0.48), w, 0.48, boxstyle="square,pad=0", facecolor="#E8EEF4", linewidth=0, zorder=2))
    for i, c in enumerate(("#FF5F57", "#FEBC2E", "#28C840")):
        ax.add_patch(Circle((x + 0.22 + i * 0.26, y + h - 0.24), 0.075, color=c, zorder=3))
    # nav icons
    ax.text(x + 1.15, y + h - 0.24, "<  >  刷新", ha="left", va="center", color="#64748B", fontsize=8, zorder=3)
    # URL bar
    ax.add_patch(FancyBboxPatch((x + 2.05, y + h - 0.38), w - 2.7, 0.28, boxstyle="round,pad=0.01,rounding_size=0.08", facecolor="white", edgecolor="#CBD5E1", linewidth=1, zorder=3))
    ax.text(x + 2.2, y + h - 0.24, url, ha="left", va="center", color="#475569", fontsize=8, zorder=4)
    # tab strip under title
    tab_y = y + h - 0.78
    ax.add_patch(FancyBboxPatch((x, tab_y), w, 0.3, boxstyle="square,pad=0", facecolor="#F1F5F9", linewidth=0, zorder=2))
    ax.add_patch(FancyBboxPatch((x + 0.15, tab_y), 1.8, 0.3, boxstyle="square,pad=0", facecolor="white", linewidth=0, zorder=3))
    ax.plot([x, x + w], [tab_y, tab_y], color="#CBD5E1", lw=0.8, zorder=4)
    ax.text(x + 1.05, tab_y + 0.15, "应用工作台", ha="center", va="center", color="#1B3A4B", fontsize=8, zorder=5)
    # app top bar
    app_top = tab_y - 0.42
    ax.add_patch(FancyBboxPatch((x, app_top), w, 0.42, boxstyle="square,pad=0", facecolor="#1B3A4B", linewidth=0, zorder=2))
    ax.text(x + 0.25, app_top + 0.21, "大模型在线应用平台", ha="left", va="center", color="white", fontsize=10, fontweight="bold", zorder=3)
    ax.text(x + w - 0.35, app_top + 0.21, "专工账号", ha="right", va="center", color="#CBD5E1", fontsize=8, zorder=3)
    cx, cy = x + 0.02, y + 0.02
    cw, ch = w - 0.04, app_top - cy
    return cx, cy, cw, ch


def mock_ui_chat(path: Path):
    fig, ax = plt.subplots(figsize=(7.8, 5.8), dpi=160)
    ax.set_xlim(0, 7.8)
    ax.set_ylim(0, 5.8)
    ax.axis("off")
    fig.patch.set_facecolor("#E2E8F0")
    cx, cy, cw, ch = _browser_frame(ax, url="https://ai-app.platform/apps/qa")
    # left nav
    ax.add_patch(FancyBboxPatch((cx, cy), 1.85, ch, boxstyle="square,pad=0", facecolor="#F8FAFC", linewidth=0, zorder=2))
    ax.plot([cx + 1.85, cx + 1.85], [cy, cy + ch], color="#E2E8F0", lw=1, zorder=3)
    ax.text(cx + 0.2, cy + ch - 0.28, "历史会话", color="#64748B", fontsize=8, zorder=3)
    for i, (t, active) in enumerate([("#3机组问答", True), ("过热器查数", False), ("规程续问", False)]):
        y = cy + ch - 0.85 - i * 0.55
        ax.add_patch(FancyBboxPatch((cx + 0.12, y), 1.6, 0.45, boxstyle="round,pad=0.02,rounding_size=0.06", facecolor="#CCFBF1" if active else "white", edgecolor="#99F6E4" if active else "#E2E8F0", linewidth=1, zorder=3))
        ax.text(cx + 0.92, y + 0.22, t, ha="center", va="center", color="#0F766E" if active else "#475569", fontsize=8, zorder=4)
    # main
    mx = cx + 1.95
    ax.text(mx + 0.1, cy + ch - 0.28, "AI问答  ·  知识检索 + 自然语言查数", color="#1B3A4B", fontsize=10, fontweight="bold", zorder=3)
    ax.text(mx + 0.1, cy + ch - 0.55, "当前话题：#3机组过热器", color="#94A3B8", fontsize=8, zorder=3)
    # user
    ax.add_patch(FancyBboxPatch((mx + 2.55, cy + ch - 1.35), 2.85, 0.55, boxstyle="round,pad=0.02,rounding_size=0.1", facecolor="#0F766E", linewidth=0, zorder=3))
    ax.text(mx + 3.95, cy + ch - 1.07, "它的超温次数上个月有多少？", ha="center", va="center", color="white", fontsize=9, zorder=4)
    # assistant card
    ax.add_patch(FancyBboxPatch((mx + 0.05, cy + ch - 3.05), 4.0, 1.5, boxstyle="round,pad=0.02,rounding_size=0.1", facecolor="white", edgecolor="#E2E8F0", linewidth=1.2, zorder=3))
    ax.add_patch(FancyBboxPatch((mx + 0.05, cy + ch - 1.75), 4.0, 0.28, boxstyle="square,pad=0", facecolor="#ECFDF5", linewidth=0, zorder=4))
    ax.text(mx + 0.2, cy + ch - 1.61, "助手  ·  台账只读查询", color="#0F766E", fontsize=8, fontweight="bold", zorder=5)
    ax.text(mx + 0.2, cy + ch - 2.15, "#3机组过热器  2026-06  超温 12 次", color="#1A2332", fontsize=10, zorder=5)
    ax.text(mx + 0.2, cy + ch - 2.5, "最高壁温 587℃    来源：运行台账 / 测点表", color="#64748B", fontsize=8, zorder=5)
    ax.text(mx + 0.2, cy + ch - 2.85, "指代「它」已解析为：#3机组过热器", color="#0F766E", fontsize=8, zorder=5)
    # chips
    ax.text(mx + 0.05, cy + 0.95, "推荐追问", color="#94A3B8", fontsize=8, zorder=3)
    for i, t in enumerate(["按受热面拆分", "对比去年同期", "相关规程条款"]):
        ax.add_patch(FancyBboxPatch((mx + 0.05 + i * 1.4, cy + 0.45), 1.3, 0.38, boxstyle="round,pad=0.02,rounding_size=0.1", facecolor="white", edgecolor="#0F766E", linewidth=1, zorder=3))
        ax.text(mx + 0.7 + i * 1.4, cy + 0.64, t, ha="center", va="center", color="#0F766E", fontsize=8, zorder=4)
    # composer
    ax.add_patch(FancyBboxPatch((mx + 0.05, cy + 0.08), 4.55, 0.32, boxstyle="round,pad=0.01,rounding_size=0.08", facecolor="#F8FAFC", edgecolor="#CBD5E1", linewidth=1, zorder=3))
    ax.text(mx + 0.2, cy + 0.24, "继续提问…", color="#94A3B8", fontsize=8, zorder=4)
    ax.add_patch(FancyBboxPatch((mx + 3.85, cy + 0.11), 0.7, 0.26, boxstyle="round,pad=0.01,rounding_size=0.06", facecolor="#0F766E", linewidth=0, zorder=4))
    ax.text(mx + 4.2, cy + 0.24, "发送", ha="center", va="center", color="white", fontsize=8, zorder=5)
    fig.tight_layout(pad=0.1)
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def mock_ui_analysis(path: Path):
    fig, ax = plt.subplots(figsize=(7.8, 5.8), dpi=160)
    ax.set_xlim(0, 7.8)
    ax.set_ylim(0, 5.8)
    ax.axis("off")
    fig.patch.set_facecolor("#E2E8F0")
    cx, cy, cw, ch = _browser_frame(ax, url="https://ai-app.platform/apps/analysis/overtemp")
    ax.text(cx + 0.25, cy + ch - 0.3, "超温分析任务", color="#1B3A4B", fontsize=11, fontweight="bold", zorder=3)
    ax.add_patch(FancyBboxPatch((cx + 5.9, cy + ch - 0.45), 1.35, 0.35, boxstyle="round,pad=0.02,rounding_size=0.08", facecolor="#0F766E", linewidth=0, zorder=3))
    ax.text(cx + 6.57, cy + ch - 0.27, "生成报告", ha="center", va="center", color="white", fontsize=9, zorder=4)
    # filter bar
    ax.add_patch(FancyBboxPatch((cx + 0.2, cy + ch - 1.05), 7.2, 0.5, boxstyle="round,pad=0.02,rounding_size=0.08", facecolor="#F8FAFC", edgecolor="#E2E8F0", linewidth=1, zorder=3))
    for i, t in enumerate(["机组  #3", "受热面  过热器", "时间  近7天", "粒度  按日"]):
        ax.text(cx + 0.55 + i * 1.75, cy + ch - 0.8, t, color="#334155", fontsize=8, zorder=4)
    # KPI
    cards = [("超温次数", "9", "次"), ("最高壁温", "591", "℃"), ("持续最长", "46", "min"), ("数据完整度", "92", "%")]
    for i, (k, v, u) in enumerate(cards):
        x = cx + 0.2 + i * 1.85
        ax.add_patch(FancyBboxPatch((x, cy + ch - 2.2), 1.7, 0.95, boxstyle="round,pad=0.02,rounding_size=0.1", facecolor="white", edgecolor="#E2E8F0", linewidth=1.2, zorder=3))
        ax.text(x + 0.15, cy + ch - 1.45, k, color="#94A3B8", fontsize=8, zorder=4)
        ax.text(x + 0.15, cy + ch - 1.85, v, color="#1B3A4B", fontsize=16, fontweight="bold", zorder=4)
        ax.text(x + 1.35, cy + ch - 1.85, u, color="#64748B", fontsize=8, zorder=4)
    # report + mini chart area
    ax.add_patch(FancyBboxPatch((cx + 0.2, cy + 0.2), 4.7, 2.35, boxstyle="round,pad=0.02,rounding_size=0.1", facecolor="white", edgecolor="#E2E8F0", linewidth=1.2, zorder=3))
    ax.text(cx + 0.4, cy + 2.3, "报告大纲（自动生成）", color="#1B3A4B", fontsize=9, fontweight="bold", zorder=4)
    for i, t in enumerate(["1  超温概况与时段分布", "2  测点/受热面对比", "3  规程对照要点", "4  初判结论与建议", "5  数据缺口说明"]):
        ax.text(cx + 0.45, cy + 1.9 - i * 0.3, t, color="#334155", fontsize=9, zorder=4)
    ax.add_patch(FancyBboxPatch((cx + 5.1, cy + 0.2), 2.3, 2.35, boxstyle="round,pad=0.02,rounding_size=0.1", facecolor="white", edgecolor="#E2E8F0", linewidth=1.2, zorder=3))
    ax.text(cx + 5.3, cy + 2.3, "近7日超温次数", color="#1B3A4B", fontsize=9, fontweight="bold", zorder=4)
    vals = [1, 2, 0, 3, 1, 1, 1]
    for i, v in enumerate(vals):
        bx = cx + 5.35 + i * 0.28
        ax.add_patch(FancyBboxPatch((bx, cy + 0.45), 0.2, 0.25 + v * 0.35, boxstyle="square,pad=0", facecolor="#0F766E", linewidth=0, zorder=4))
    fig.tight_layout(pad=0.1)
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def mock_ui_vision(path: Path):
    fig, ax = plt.subplots(figsize=(7.8, 5.8), dpi=160)
    ax.set_xlim(0, 7.8)
    ax.set_ylim(0, 5.8)
    ax.axis("off")
    fig.patch.set_facecolor("#E2E8F0")
    cx, cy, cw, ch = _browser_frame(ax, url="https://ai-app.platform/apps/vision-diagnose")
    ax.text(cx + 0.25, cy + ch - 0.28, "看图诊断  ·  缺陷识别", color="#1B3A4B", fontsize=11, fontweight="bold", zorder=3)
    # photo panel
    img_l, img_b, img_w, img_h = cx + 0.2, cy + 1.15, 3.55, 3.35
    ax.add_patch(FancyBboxPatch((img_l - 0.05, img_b - 0.05), img_w + 0.1, img_h + 0.1, boxstyle="round,pad=0.01,rounding_size=0.08", facecolor="#0F172A", linewidth=0, zorder=2))
    if PHOTO_DEFECT.exists():
        ax.imshow(_load_rgb(PHOTO_DEFECT), extent=(img_l, img_l + img_w, img_b, img_b + img_h), aspect="auto", zorder=3)
    # detection overlay (approx on chalk circle area)
    ax.add_patch(FancyBboxPatch((img_l + 0.85, img_b + 1.15), 1.7, 1.35, boxstyle="square,pad=0", facecolor="none", edgecolor="#FBBF24", linewidth=2.2, zorder=5))
    ax.add_patch(FancyBboxPatch((img_l + 0.85, img_b + 2.35), 1.15, 0.28, boxstyle="square,pad=0", facecolor="#FBBF24", linewidth=0, zorder=6))
    ax.text(img_l + 1.42, img_b + 2.49, "裂纹 0.91", ha="center", va="center", color="#1A2332", fontsize=8, fontweight="bold", zorder=7)
    ax.text(img_l + 0.1, img_b + 0.12, "现场照片（已标注疑似缺陷区）", color="white", fontsize=7, zorder=6)
    # right panel
    rx = cx + 4.0
    ax.add_patch(FancyBboxPatch((rx, cy + 1.15), 3.35, 3.35, boxstyle="round,pad=0.02,rounding_size=0.1", facecolor="white", edgecolor="#E2E8F0", linewidth=1.2, zorder=3))
    ax.text(rx + 0.2, cy + 4.2, "研判结果", color="#1B3A4B", fontsize=10, fontweight="bold", zorder=4)
    ax.add_patch(FancyBboxPatch((rx + 0.2, cy + 3.6), 1.2, 0.35, boxstyle="round,pad=0.02,rounding_size=0.08", facecolor="#FEF3C7", edgecolor="#F59E0B", linewidth=1, zorder=4))
    ax.text(rx + 0.8, cy + 3.77, "风险：中", ha="center", va="center", color="#B45309", fontsize=9, fontweight="bold", zorder=5)
    for i, t in enumerate(["视觉：管壁锈蚀区可见横向裂纹", "历史：近30天同类缺陷 2 条", "规程：建议复核测厚并评估减薄", "处置：隔离观察 · 安排专项检查"]):
        ax.text(rx + 0.2, cy + 3.2 - i * 0.4, "·  " + t, color="#334155", fontsize=8, zorder=4)
    ax.add_patch(FancyBboxPatch((rx + 0.2, cy + 1.35), 2.95, 0.55, boxstyle="round,pad=0.02,rounding_size=0.08", facecolor="#ECFDF5", linewidth=0, zorder=4))
    ax.text(rx + 1.67, cy + 1.62, "已联动：历史台账 + 企业知识库", ha="center", va="center", color="#0F766E", fontsize=8, zorder=5)
    # confirm bar
    ax.add_patch(FancyBboxPatch((cx + 0.2, cy + 0.2), 7.2, 0.8, boxstyle="round,pad=0.02,rounding_size=0.1", facecolor="#FFF7ED", edgecolor="#FDBA74", linewidth=1, zorder=3))
    ax.text(cx + 0.4, cy + 0.75, "人机确认：请确认检查对象与范围后继续", color="#B45309", fontsize=9, fontweight="bold", zorder=4)
    ax.text(cx + 0.4, cy + 0.4, "对象：受热面管排  ·  区域：标注裂纹区（白圈附近）", color="#78716C", fontsize=8, zorder=4)
    ax.add_patch(FancyBboxPatch((cx + 5.55, cy + 0.35), 0.85, 0.4, boxstyle="round,pad=0.02,rounding_size=0.08", facecolor="#0F766E", linewidth=0, zorder=4))
    ax.text(cx + 5.97, cy + 0.55, "确认", ha="center", va="center", color="white", fontsize=9, zorder=5)
    ax.add_patch(FancyBboxPatch((cx + 6.55, cy + 0.35), 0.7, 0.4, boxstyle="round,pad=0.02,rounding_size=0.08", facecolor="white", edgecolor="#A8A29E", linewidth=1, zorder=4))
    ax.text(cx + 6.9, cy + 0.55, "修改", ha="center", va="center", color="#57534E", fontsize=9, zorder=5)
    fig.tight_layout(pad=0.1)
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def mock_ui_doc(path: Path):
    fig, ax = plt.subplots(figsize=(7.8, 5.8), dpi=160)
    ax.set_xlim(0, 7.8)
    ax.set_ylim(0, 5.8)
    ax.axis("off")
    fig.patch.set_facecolor("#E2E8F0")
    cx, cy, cw, ch = _browser_frame(ax, url="https://ai-app.platform/apps/doc-extract")
    ax.text(cx + 0.25, cy + ch - 0.28, "检修报告结构化提取", color="#1B3A4B", fontsize=11, fontweight="bold", zorder=3)
    ax.text(cx + 0.25, cy + ch - 0.55, "任务状态：抽取完成 · 待复核", color="#0F766E", fontsize=8, zorder=3)
    # left document viewer
    ax.add_patch(FancyBboxPatch((cx + 0.2, cy + 0.2), 3.35, 3.85, boxstyle="round,pad=0.02,rounding_size=0.1", facecolor="#F8FAFC", edgecolor="#E2E8F0", linewidth=1.2, zorder=3))
    ax.text(cx + 0.4, cy + 3.8, "原文预览  ·  检修报告.pdf", color="#64748B", fontsize=8, zorder=4)
    lines = [
        "三、缺陷检查记录",
        "检查部位：#3机组过热器右侧管排",
        "发现磨损减薄，局部可见锈蚀与裂纹。",
        "实测壁厚约 3.2mm，建议复核测厚。",
        "发现日期：2026-06-18",
        "四、处理意见（略）",
    ]
    for i, t in enumerate(lines):
        y = cy + 3.35 - i * 0.42
        bg = "#FEF3C7" if i in (1, 2, 3) else "white"
        ax.add_patch(FancyBboxPatch((cx + 0.35, y), 3.05, 0.35, boxstyle="round,pad=0.01,rounding_size=0.04", facecolor=bg, edgecolor="#FDE68A" if i in (1, 2, 3) else "#F1F5F9", linewidth=1, zorder=4))
        ax.text(cx + 0.45, y + 0.17, t, color="#92400E" if i in (1, 2, 3) else "#475569", fontsize=7.5, zorder=5)
    # right form
    ax.add_patch(FancyBboxPatch((cx + 3.8, cy + 0.2), 3.55, 3.85, boxstyle="round,pad=0.02,rounding_size=0.1", facecolor="white", edgecolor="#99F6E4", linewidth=1.5, zorder=3))
    ax.text(cx + 4.0, cy + 3.8, "标准缺陷记录（可编辑）", color="#0F766E", fontsize=9, fontweight="bold", zorder=4)
    fields = [("机组", "#3"), ("位置", "过热器右侧管排"), ("缺陷类型", "磨损减薄/裂纹"), ("壁厚(mm)", "3.2"), ("发现日期", "2026-06-18"), ("状态", "待复核入库")]
    for i, (k, v) in enumerate(fields):
        y = cy + 3.25 - i * 0.48
        ax.text(cx + 4.0, y + 0.12, k, color="#94A3B8", fontsize=7.5, zorder=4)
        ax.add_patch(FancyBboxPatch((cx + 5.1, y), 2.0, 0.34, boxstyle="round,pad=0.01,rounding_size=0.05", facecolor="#F0FDFA", edgecolor="#99F6E4", linewidth=1, zorder=4))
        ax.text(cx + 6.1, y + 0.17, v, ha="center", va="center", color="#134E4A", fontsize=8, zorder=5)
    ax.add_patch(FancyBboxPatch((cx + 4.0, cy + 0.35), 1.4, 0.38, boxstyle="round,pad=0.02,rounding_size=0.08", facecolor="#0F766E", linewidth=0, zorder=4))
    ax.text(cx + 4.7, cy + 0.54, "确认入库", ha="center", va="center", color="white", fontsize=9, zorder=5)
    fig.tight_layout(pad=0.1)
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def mock_ui_rootcause(path: Path):
    fig, ax = plt.subplots(figsize=(7.8, 5.8), dpi=160)
    ax.set_xlim(0, 7.8)
    ax.set_ylim(0, 5.8)
    ax.axis("off")
    fig.patch.set_facecolor("#E2E8F0")
    cx, cy, cw, ch = _browser_frame(ax, url="https://ai-app.platform/apps/root-cause")
    ax.text(cx + 0.25, cy + ch - 0.28, "泄爆溯源分析", color="#1B3A4B", fontsize=11, fontweight="bold", zorder=3)
    ax.text(cx + 0.25, cy + ch - 0.55, "模板：分方向 · 分层次（直接原因→中期因素→根因）", color="#64748B", fontsize=8, zorder=3)
    dirs = [
        ("材料方向", "材质/老化", "焊接/热处理", "供货批次"),
        ("运行方向", "超温超压", "水质工况", "启停冲击"),
        ("检修方向", "测厚漏检", "防磨缺失", "工艺执行"),
        ("设计方向", "结构应力", "流速冲刷", "选材裕度"),
    ]
    for i, (title, a, b, c) in enumerate(dirs):
        x = cx + 0.2 + (i % 2) * 3.7
        y = cy + ch - 2.35 - (i // 2) * 1.55
        color = "#0F766E" if i % 2 == 0 else "#1B3A4B"
        ax.add_patch(FancyBboxPatch((x, y), 3.5, 1.4, boxstyle="round,pad=0.02,rounding_size=0.1", facecolor="white", edgecolor="#E2E8F0", linewidth=1.2, zorder=3))
        ax.add_patch(FancyBboxPatch((x, y + 1.05), 3.5, 0.35, boxstyle="square,pad=0", facecolor=color, linewidth=0, zorder=4))
        ax.text(x + 1.75, y + 1.22, title, ha="center", va="center", color="white", fontsize=9, fontweight="bold", zorder=5)
        ax.text(x + 0.2, y + 0.7, f"直接：{a}", color="#334155", fontsize=8, zorder=4)
        ax.text(x + 0.2, y + 0.42, f"中期：{b}", color="#334155", fontsize=8, zorder=4)
        ax.text(x + 0.2, y + 0.14, f"根因：{c}  ·  已挂接规程/案例", color="#0F766E", fontsize=8, zorder=4)
    ax.add_patch(FancyBboxPatch((cx + 0.2, cy + 0.2), 7.2, 0.5, boxstyle="round,pad=0.02,rounding_size=0.08", facecolor="#ECFDF5", edgecolor="#99F6E4", linewidth=1, zorder=3))
    ax.text(cx + 3.8, cy + 0.45, "输出草稿可供调查讨论  ·  支持导出结构化报告", ha="center", va="center", color="#0F766E", fontsize=9, zorder=4)
    fig.tight_layout(pad=0.1)
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def prepare_risk_photo(src: Path | None = None, dest: Path | None = None) -> Path:
    """Mosaic sensitive text and draw detection boxes on risk site photo."""
    from PIL import Image, ImageDraw

    dest = dest or PHOTO_RISK
    raw = ASSETS / "photo_risk_site_raw.png"
    if src and Path(src).exists():
        Image.open(src).convert("RGB").save(raw)
    if not raw.exists() and dest.exists():
        # fallback: already prepared
        return dest
    if not raw.exists():
        return dest

    im = Image.open(raw).convert("RGB")

    def mosaic(box, block=6):
        x1, y1, x2, y2 = [int(v) for v in box]
        x1, y1 = max(0, x1), max(0, y1)
        x2, y2 = min(im.size[0], x2), min(im.size[1], y2)
        region = im.crop((x1, y1, x2, y2))
        rw, rh = region.size
        if rw < 2 or rh < 2:
            return
        small = region.resize((max(1, rw // block), max(1, rh // block)), Image.BILINEAR)
        im.paste(small.resize((rw, rh), Image.NEAREST), (x1, y1))

    mosaic((0, 0, 450, 70))
    mosaic((560, 485, 1024, 576))

    draw = ImageDraw.Draw(im)

    def box_label(xy, label, color, text_color=(255, 255, 255)):
        x1, y1, x2, y2 = xy
        draw.rectangle([x1, y1, x2, y2], outline=color, width=3)
        tw = max(70, 8 * len(label) + 8)
        th = 20
        ly = max(0, y1 - th - 2)
        draw.rectangle([x1, ly, x1 + tw, ly + th], fill=color)
        draw.text((x1 + 4, ly + 3), label, fill=text_color)

    box_label((450, 110, 740, 410), "吊车 0.93", (56, 189, 248), (15, 23, 42))
    box_label((320, 390, 385, 485), "人员 0.92", (248, 113, 113))
    box_label((355, 405, 415, 495), "人员 0.90", (248, 113, 113))
    box_label((305, 145, 365, 215), "人员 0.89", (248, 113, 113))
    box_label((230, 185, 285, 255), "人员 0.87", (248, 113, 113))
    box_label((500, 170, 560, 260), "人员 0.91", (248, 113, 113))
    box_label((650, 120, 705, 200), "人员 0.88", (248, 113, 113))
    im.save(dest)
    return dest


def mock_ui_risk(path: Path):
    fig, ax = plt.subplots(figsize=(7.8, 5.8), dpi=160)
    ax.set_xlim(0, 7.8)
    ax.set_ylim(0, 5.8)
    ax.axis("off")
    fig.patch.set_facecolor("#E2E8F0")
    cx, cy, cw, ch = _browser_frame(ax, url="https://ai-app.platform/apps/risk-behavior")
    ax.text(cx + 0.25, cy + ch - 0.28, "作业风险行为研判", color="#1B3A4B", fontsize=11, fontweight="bold", zorder=3)
    ax.text(cx + 0.25, cy + ch - 0.52, "作业现场监控画面  ·  目标检测 + 二次研判", color="#94A3B8", fontsize=8, zorder=3)
    # live frame with annotated photo (text mosaiced, boxes baked in)
    prepare_risk_photo()
    img_l, img_b, img_w, img_h = cx + 0.2, cy + 1.25, 4.35, 2.95
    ax.add_patch(FancyBboxPatch((img_l - 0.04, img_b - 0.04), img_w + 0.08, img_h + 0.08, boxstyle="round,pad=0.01,rounding_size=0.06", facecolor="#0F172A", linewidth=0, zorder=2))
    if PHOTO_RISK.exists():
        ax.imshow(_load_rgb(PHOTO_RISK), extent=(img_l, img_l + img_w, img_b, img_b + img_h), aspect="auto", zorder=3)
    ax.text(img_l + 0.1, img_b + 0.1, "小模型目标检测叠加（敏感信息已脱敏）", color="white", fontsize=7, zorder=6)
    # right judgment
    rx = cx + 4.75
    ax.add_patch(FancyBboxPatch((rx, cy + 1.25), 2.6, 2.95, boxstyle="round,pad=0.02,rounding_size=0.1", facecolor="white", edgecolor="#E2E8F0", linewidth=1.2, zorder=3))
    ax.text(rx + 0.15, cy + 3.95, "大模型二次研判", color="#1B3A4B", fontsize=9, fontweight="bold", zorder=4)
    ax.add_patch(FancyBboxPatch((rx + 0.15, cy + 3.35), 2.3, 0.4, boxstyle="round,pad=0.02,rounding_size=0.08", facecolor="#FEE2E2", edgecolor="#EF4444", linewidth=1, zorder=4))
    ax.text(rx + 1.3, cy + 3.55, "确认风险 · 高", ha="center", va="center", color="#DC2626", fontsize=10, fontweight="bold", zorder=5)
    for i, t in enumerate(["行为：吊装作业区人员停留", "语义：交叉作业风险", "分级：高 · 建议立即处置", "状态：已回写预警链路"]):
        ax.text(rx + 0.15, cy + 2.95 - i * 0.35, "· " + t, color="#334155", fontsize=8, zorder=4)
    # pipeline
    steps = ["小模型检出", "上送画面", "大模型研判", "自动分级"]
    for i, t in enumerate(steps):
        x = cx + 0.25 + i * 1.85
        ax.add_patch(FancyBboxPatch((x, cy + 0.25), 1.65, 0.7, boxstyle="round,pad=0.02,rounding_size=0.1", facecolor="#0F766E" if i % 2 == 0 else "#1B3A4B", linewidth=0, zorder=3))
        ax.text(x + 0.82, cy + 0.6, t, ha="center", va="center", color="white", fontsize=9, fontweight="bold", zorder=4)
        if i < 3:
            ax.annotate("", xy=(x + 1.78, cy + 0.6), xytext=(x + 1.68, cy + 0.6), arrowprops=dict(arrowstyle="-|>", color="#64748B", lw=1.2), zorder=4)
    fig.tight_layout(pad=0.1)
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def build_charts():
    paths = {
        "overview": ASSETS / "overview.png",
        "industry": ASSETS / "industry.png",
        "scenes": ASSETS / "scenes.png",
        "solmap": ASSETS / "solmap.png",
        "flows": ASSETS / "flows.png",
        "d1": ASSETS / "detail_qa.png",
        "d2": ASSETS / "detail_analysis.png",
        "d3": ASSETS / "detail_vision.png",
        "d4": ASSETS / "detail_doc.png",
        "d5": ASSETS / "detail_risk.png",
    }
    chart_overview(paths["overview"])
    chart_industry(paths["industry"])
    chart_scenes(paths["scenes"])
    chart_solution_map(paths["solmap"])
    chart_solution_flows(paths["flows"])
    chart_detail_pipeline(
        paths["d1"],
        "解决路径",
        ["用户提问", "意图分流", "知识检索 / 自然语言查数", "流式回答 / 澄清追问"],
        "#0F766E",
    )
    chart_detail_pipeline(
        paths["d2"],
        "解决路径",
        ["分析需求", "取数计划", "并行查数 + 知识检索", "质量检查", "结构化报告"],
        "#1B3A4B",
    )
    chart_detail_pipeline(
        paths["d3"],
        "解决路径",
        ["现场图像 + 描述", "视觉研判 ‖ 历史取数 ‖ 规范检索", "合成研判建议", "必要时人工确认续跑"],
        "#0F766E",
    )
    chart_detail_pipeline(
        paths["d4"],
        "解决路径",
        ["业务文档上传", "解析 / 分块", "多阶段抽取", "标准记录输出", "复核入库"],
        "#1B3A4B",
    )
    chart_detail_pipeline(
        paths["d5"],
        "解决路径",
        ["小模型目标检测", "上送图像 / 目标框", "多模态大模型二次研判", "确认风险并自动分级", "回写预警"],
        "#0F766E",
    )
    paths.update(
        {
            "ui_chat": ASSETS / "ui_chat.png",
            "ui_analysis": ASSETS / "ui_analysis.png",
            "ui_vision": ASSETS / "ui_vision.png",
            "ui_doc": ASSETS / "ui_doc.png",
            "ui_rootcause": ASSETS / "ui_rootcause.png",
            "ui_risk": ASSETS / "ui_risk.png",
        }
    )
    mock_ui_chat(paths["ui_chat"])
    mock_ui_analysis(paths["ui_analysis"])
    mock_ui_vision(paths["ui_vision"])
    mock_ui_doc(paths["ui_doc"])
    mock_ui_rootcause(paths["ui_rootcause"])
    mock_ui_risk(paths["ui_risk"])
    return paths


def add_scene_demo_slide(prs, blank, page, total, title, mapping, summary, app_points, chart_path):
    s = prs.slides.add_slide(blank)
    page_header(s, title, mapping, section="PART 02  ·  场景效果示意")
    # left mockup
    rounded_rect(s, Inches(0.45), Inches(1.35), Inches(7.55), Inches(5.45), fill=C["white"], line=C["line"])
    s.shapes.add_picture(str(chart_path), Inches(0.55), Inches(1.45), width=Inches(7.35))
    # right text
    rounded_rect(s, Inches(8.2), Inches(1.35), Inches(4.55), Inches(5.45), fill=C["white"], line=C["line"])
    accent_bar(s, Inches(8.2), Inches(1.35), width=Inches(0.1), height=Inches(5.45), color=C["teal"])
    add_textbox(s, Inches(8.5), Inches(1.55), Inches(4.0), Inches(0.35), "场景说明", size=14, bold=True, color=C["navy"])
    add_textbox(s, Inches(8.5), Inches(1.95), Inches(4.0), Inches(1.35), summary, size=12, color=C["ink"])
    add_textbox(s, Inches(8.5), Inches(3.4), Inches(4.0), Inches(0.35), "应用场景", size=14, bold=True, color=C["navy"])
    tb = s.shapes.add_textbox(Inches(8.5), Inches(3.85), Inches(4.0), Inches(2.6))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for i, line in enumerate(app_points[:2], 1):
        txt = f"{i}.  {line}"
        if first:
            set_run(tf.paragraphs[0].add_run(), txt, size=12, color=C["muted"])
            first = False
        else:
            add_para(tf, txt, size=12, color=C["muted"], space_before=14)
    footer(s, page, total)
    return s


def build_ppt(charts: dict):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 22

    # 1 Cover
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, C["navy"])
    deco = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), Inches(7.5))
    fill_shape(deco, C["teal"])
    add_textbox(s, Inches(1.1), Inches(2.0), Inches(11), Inches(0.4), "企业级 AI 应用解决方案", size=16, color="99F6E4", bold=True)
    add_textbox(s, Inches(1.1), Inches(2.55), Inches(11), Inches(1.0), "成果积累与应用场景示例", size=36, bold=True, color=C["white"])
    add_textbox(s, Inches(1.1), Inches(3.7), Inches(10.5), Inches(0.7), "行业场景解决方案  ·  典型应用示例", size=18, color="CBD5E1")
    add_textbox(s, Inches(1.1), Inches(5.6), Inches(10), Inches(0.4), "能源电力工业场景  |  大模型在线应用平台", size=13, color="94A3B8")
    footer(s, 1, total, light=False)

    # 2 Agenda
    s = prs.slides.add_slide(blank)
    page_header(s, "汇报结构", "行业方案成果 + 应用场景示例")
    agenda = [
        ("01", "定位与主张", "我们交付什么"),
        ("02", "行业场景方案", "总览 + 要点 + 详解"),
        ("03", "应用场景示例", "项目落地实例"),
    ]
    for i, (num, title, desc) in enumerate(agenda):
        left = Inches(1.55 + i * 3.7)
        rounded_rect(s, left, Inches(2.0), Inches(3.4), Inches(3.5), fill=C["white"], line=C["line"])
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.25), Inches(2.4), Inches(0.9), Inches(0.9))
        fill_shape(circle, C["teal"] if i % 2 == 0 else C["navy"])
        add_textbox(s, left + Inches(1.25), Inches(2.58), Inches(0.9), Inches(0.55), num, size=18, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.2), Inches(3.55), Inches(3.0), Inches(0.5), title, size=15, bold=True, color=C["navy"], align=PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.25), Inches(4.2), Inches(2.9), Inches(0.8), desc, size=12, color=C["muted"], align=PP_ALIGN.CENTER)
    footer(s, 2, total)

    # 3 Positioning
    s = prs.slides.add_slide(blank)
    page_header(s, "定位与核心主张", "行业可交付 · 场景可落地 · 生产可用")
    rounded_rect(s, Inches(0.7), Inches(1.55), Inches(11.9), Inches(1.25), fill=C["white"], line=C["line"])
    add_textbox(
        s,
        Inches(1.0),
        Inches(1.8),
        Inches(11.3),
        Inches(0.8),
        "面向能源及同类工业场景，沉淀可跨行业复用的解决方案包，并交付可落地的业务应用。",
        size=15,
        color=C["ink"],
    )
    vals = [
        ("行业方案可复制", "问答、分析、现场研判、文档整理\n按场景打包，换客户主要接知识与数据"),
        ("场景可落地", "通用方案已有项目实例\n可演示、可推广"),
        ("双轮驱动", "企业知识\n+ 业务数据协同"),
        ("生产可用", "只读查数 · 知识隔离\n关键环节可人工确认"),
    ]
    for i, (t, b) in enumerate(vals):
        left = Inches(0.7 + i * 3.1)
        rounded_rect(s, left, Inches(3.15), Inches(2.95), Inches(3.0), fill=C["white"], line=C["line"])
        accent_bar(s, left, Inches(3.15), width=Inches(0.1), height=Inches(3.0), color=C["teal"] if i % 2 == 0 else C["navy"])
        add_textbox(s, left + Inches(0.25), Inches(3.4), Inches(2.5), Inches(0.5), t, size=14, bold=True, color=C["navy"])
        add_textbox(s, left + Inches(0.25), Inches(4.1), Inches(2.5), Inches(1.6), b, size=12, color=C["muted"])
    footer(s, 3, total)

    # 4 Overview chart
    s = prs.slides.add_slide(blank)
    page_header(s, "成果总览", "行业方案包 + 已落地场景", section="PART 01  ·  成果积累")
    s.shapes.add_picture(str(charts["overview"]), Inches(0.5), Inches(1.5), width=Inches(12.3))
    footer(s, 4, total)

    # 5 Industry overview
    s = prs.slides.add_slide(blank)
    page_header(s, "行业场景解决方案", "回答「业务怎么用、解决谁的问题」", section="PART 01  ·  成果积累")
    s.shapes.add_picture(str(charts["industry"]), Inches(0.55), Inches(1.5), width=Inches(12.2))
    footer(s, 5, total)

    # 6 Industry details
    s = prs.slides.add_slide(blank)
    page_header(s, "行业方案要点", "通用命名：问题 → 价值（跨行业可套用）", section="PART 01  ·  成果积累")
    industry = [
        ("知识问答与数据查询一体化", "资料分散、台账难查", "统一入口问知识、查数据，支持多轮续问"),
        ("多源业务综合分析", "分析拼数慢、格式乱", "自动取数生成结构化分析报告"),
        ("多模态现场研判", "拍照后仍要回办公室查", "图像联历史与规范，辅助现场处置"),
        ("业务文档结构化整理", "人工录入慢、字段乱", "文档自动整理为可入库标准记录"),
        ("作业风险行为检测与研判", "复杂现场误报高、难懂行为语义", "小模型检出后，多模态大模型二次确认并分级"),
    ]
    # top row 3
    for i, (t, p, v) in enumerate(industry[:3]):
        left = Inches(0.55 + i * 4.2)
        top = Inches(1.45)
        rounded_rect(s, left, top, Inches(4.0), Inches(2.4), fill=C["white"], line=C["line"])
        accent_bar(s, left, top, width=Inches(0.1), height=Inches(2.4), color=C["teal"] if i % 2 == 0 else C["navy"])
        add_textbox(s, left + Inches(0.25), top + Inches(0.2), Inches(3.5), Inches(0.45), t, size=14, bold=True, color=C["navy"])
        add_textbox(s, left + Inches(0.25), top + Inches(0.8), Inches(3.5), Inches(0.4), "问题：" + p, size=11, color=C["amber"])
        add_textbox(s, left + Inches(0.25), top + Inches(1.35), Inches(3.5), Inches(0.85), "价值：" + v, size=12, color=C["muted"])
    # bottom row 2 centered — slightly shorter to keep footer clear
    for i, (t, p, v) in enumerate(industry[3:]):
        left = Inches(2.15 + i * 4.6)
        top = Inches(4.15)
        rounded_rect(s, left, top, Inches(4.35), Inches(2.2), fill=C["white"], line=C["line"])
        accent_bar(s, left, top, width=Inches(0.1), height=Inches(2.2), color=C["navy"] if i % 2 == 0 else C["teal"])
        add_textbox(s, left + Inches(0.25), top + Inches(0.18), Inches(3.9), Inches(0.4), t, size=14, bold=True, color=C["navy"])
        add_textbox(s, left + Inches(0.25), top + Inches(0.7), Inches(3.9), Inches(0.35), "问题：" + p, size=11, color=C["amber"])
        add_textbox(s, left + Inches(0.25), top + Inches(1.15), Inches(3.9), Inches(0.85), "价值：" + v, size=12, color=C["muted"])
    footer(s, 6, total)

    # 7-11 One scenario per page
    add_detail_slide(
        prs,
        blank,
        7,
        total,
        "详解①  知识问答与业务数据查询一体化",
        [
            "· 制度/规程/案例分散，一线翻册、找人成本高",
            "· 台账统计常依赖专人写查询，响应慢",
            "· 多轮追问「它 / 上次那个 / 第几点」易答偏",
        ],
        [
            "· 统一入口自动分流：问知识 / 查数据",
            "· 知识库检索 + 自然语言只读查数",
            "· 支持指代消解、条目续问与推荐追问",
        ],
        [
            "· AI问答 · 知识库运营 · 安全只读查数",
            "· 多轮会话 · 流式交互",
        ],
        "AI问答（规程/案例 + 台账查数）",
        charts["d1"],
        C["teal"],
    )
    add_detail_slide(
        prs,
        blank,
        8,
        total,
        "详解②  多源业务综合分析",
        [
            "· 异常/策略/事故研判需跨系统拼数，周期长",
            "· 报告格式因人而异，交接复核成本高",
            "· 数据不全时易硬编，缺少缺口说明",
        ],
        [
            "· 按专题配置取数计划与报告模板",
            "· 多维数据与知识并行获取 + 质量门控",
            "· 输出结构化报告（含溯源类模板）",
        ],
        [
            "· 可配置专题分析服务 · 并行取数",
            "· 质量门控 · 结构化报告生成",
        ],
        "超温 / 检修策略 / 四管健康 / 泄爆溯源",
        charts["d2"],
        C["navy"],
    )
    add_detail_slide(
        prs,
        blank,
        9,
        total,
        "详解③  多模态现场研判",
        [
            "· 现场拍照后仍要回办公室查历史与规范",
            "· 图像、历史数据、规程知识彼此割裂",
            "· 对象/范围不清时全自动硬推易误判",
        ],
        [
            "· 视觉研判 ‖ 历史取数 ‖ 规范检索并行",
            "· 合成研判与处置建议",
            "· 关键处可人工确认后再续跑",
        ],
        [
            "· 图片预处理 · 多模态研判 · 历史联动",
            "· 知识检索 · 人机确认续跑",
        ],
        "看图诊断（缺陷识别 / 泄爆辅助）",
        charts["d3"],
        C["teal"],
    )
    add_detail_slide(
        prs,
        blank,
        10,
        total,
        "详解④  业务文档结构化整理",
        [
            "· Word/PDF 报告人工录入耗时长",
            "· 字段口径不统一，难入库、难统计",
            "· 扫描件与复杂表格抄录易错",
        ],
        [
            "· 长文档解析分块 → 多阶段智能抽取",
            "· 输出标准业务记录（支持扫描件/复杂表）",
            "· 结果可复核后入库",
        ],
        [
            "· 异步文档处理 · 分块并行抽取",
            "· 结构化校验输出 · 任务状态与入库衔接",
        ],
        "检修报告结构化提取",
        charts["d4"],
        C["navy"],
    )
    add_detail_slide(
        prs,
        blank,
        11,
        total,
        "详解⑤  作业风险行为检测与研判",
        [
            "· 专业小模型易误报，告警噪音大",
            "· 人—设备—环境关系与行为语义难判",
            "· 风险分级高度依赖人工复核",
        ],
        [
            "· 小模型先做目标检测并上送疑似画面",
            "· 多模态大模型二次风险行为检测与研判",
            "· 确认真伪并自动分级，回写预警链路",
        ],
        [
            "· 小模型接入 · 疑似风险上送",
            "· 大模型二次研判 · 风险自动分级",
        ],
        "安全生产准实时风险研判（大小模型协同）",
        charts["d5"],
        C["teal"],
    )

    # 12 Scenario→Solution→Delivery map
    s = prs.slides.add_slide(blank)
    page_header(s, "行业场景痛点 → 方案 → 交付", "一眼看清我们针对各类场景能提供什么", section="PART 01  ·  行业方案详解")
    s.shapes.add_picture(str(charts["solmap"]), Inches(0.45), Inches(1.35), width=Inches(12.4))
    footer(s, 12, total)

    # 13 Divider
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, C["teal"])
    deco = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), Inches(7.5))
    fill_shape(deco, C["navy"])
    add_textbox(s, Inches(1.1), Inches(2.6), Inches(11), Inches(0.4), "PART 02", size=16, bold=True, color="99F6E4")
    add_textbox(s, Inches(1.1), Inches(3.15), Inches(11), Inches(0.8), "应用场景示例", size=36, bold=True, color=C["white"])
    add_textbox(s, Inches(1.1), Inches(4.15), Inches(10), Inches(0.5), "当前项目已落地的典型业务场景与效果描述", size=15, color="CCFBF1")
    footer(s, 13, total, light=False)

    # 14 Scene map
    s = prs.slides.add_slide(blank)
    page_header(s, "场景全景", "按日常问询 / 分析研判 / 现场与文档归集", section="PART 02  ·  应用场景示例")
    s.shapes.add_picture(str(charts["scenes"]), Inches(0.55), Inches(1.45), width=Inches(12.2))
    footer(s, 14, total)

    # 15 Scenario cards
    s = prs.slides.add_slide(blank)
    page_header(s, "典型场景与效果（项目实例）", "通用方案在电力项目中的落地示例；后页为各场景前端效果示意", section="PART 02  ·  应用场景示例")
    scenes = [
        ("AI问答", "对应：知识问答与数据查询一体化", "少翻册、少找人；支持「它 / 第几点」续问。"),
        ("超温 / 专项分析", "对应：多源业务综合分析", "缩短初判与起稿时间；结构统一、缺口可说明。"),
        ("看图诊断", "对应：多模态现场研判", "现场辅助研判；范围不清可人工确认。"),
        ("检修报告提取", "对应：业务文档结构化整理", "录入由抄录转为复核，便于入库。"),
        ("泄爆溯源分析", "对应：多源业务综合分析", "事故调查有规范草稿起步，便于讨论复核。"),
        ("风险行为研判", "对应：作业风险行为检测与研判", "降低误报；理解作业语义；支持风险自动分级。"),
    ]
    for i, (t, sub, eff) in enumerate(scenes):
        r, c = divmod(i, 3)
        left = Inches(0.65 + c * 4.15)
        top = Inches(1.5 + r * 2.6)
        rounded_rect(s, left, top, Inches(3.95), Inches(2.35), fill=C["white"], line=C["line"])
        accent_bar(s, left, top, width=Inches(0.1), height=Inches(2.35), color=C["teal"] if i % 2 == 0 else C["navy"])
        add_textbox(s, left + Inches(0.25), top + Inches(0.2), Inches(3.5), Inches(0.35), t, size=15, bold=True, color=C["navy"])
        add_textbox(s, left + Inches(0.25), top + Inches(0.6), Inches(3.5), Inches(0.4), sub, size=11, color=C["teal"])
        add_textbox(s, left + Inches(0.25), top + Inches(1.15), Inches(3.5), Inches(1.0), eff, size=12, color=C["muted"])
    footer(s, 15, total)

    # 16-21 Per-scenario UI mock + narrative
    add_scene_demo_slide(
        prs,
        blank,
        16,
        total,
        "示例 · AI问答",
        "对应方案：知识问答与业务数据查询一体化",
        "统一对话入口承接规程/案例问答与台账查数；指代续问与推荐追问降低重复描述成本。",
        [
            "值班/一线查规程、案例，少翻册、少找人，7×24 可用",
            "运行台账统计用自然语言提问，不必找人写查询",
        ],
        charts["ui_chat"],
    )
    add_scene_demo_slide(
        prs,
        blank,
        17,
        total,
        "示例 · 超温 / 专项分析",
        "对应方案：多源业务综合分析",
        "按机组、受热面、时间等条件自动取数，生成结构化分析报告，缩短专工初判与起稿时间。",
        [
            "超温异常发生后，专工快速拿到结构化初判报告",
            "报告口径统一，便于交接复核与管理决策",
        ],
        charts["ui_analysis"],
    )
    add_scene_demo_slide(
        prs,
        blank,
        18,
        total,
        "示例 · 看图诊断",
        "对应方案：多模态现场研判",
        "上传现场照片后，系统并行完成视觉研判、历史联动与规程对照，必要时弹出人机确认。",
        [
            "巡检/停机现场拍照即可辅助研判，减少回办公室查资料",
            "历史缺陷与规程建议一并给出，支撑现场处置决策",
        ],
        charts["ui_vision"],
    )
    add_scene_demo_slide(
        prs,
        blank,
        19,
        total,
        "示例 · 检修报告提取",
        "对应方案：业务文档结构化整理",
        "Word/PDF 检修报告自动抽取为标准缺陷字段，录入由「逐页抄录」转为「上传后复核」。",
        [
            "检修文档岗从抄录改为复核，显著缩短入库周期",
            "缺陷字段口径统一，便于后续统计与专项分析",
        ],
        charts["ui_doc"],
    )
    add_scene_demo_slide(
        prs,
        blank,
        20,
        total,
        "示例 · 泄爆溯源分析",
        "对应方案：多源业务综合分析",
        "面向泄漏/爆管，按材料、运行、检修、设计等方向输出「直接原因→中期因素→根因」草稿。",
        [
            "事故调查有规范溯源草稿起步，减少遗漏排查方向",
            "防控建议可讨论复核，缩短调查起稿与汇报准备时间",
        ],
        charts["ui_rootcause"],
    )
    add_scene_demo_slide(
        prs,
        blank,
        21,
        total,
        "示例 · 风险行为研判",
        "对应方案：作业风险行为检测与研判",
        "小模型检出人/设备目标后，多模态大模型做二次风险行为确认与自动分级，并回写预警。",
        [
            "安监值守降低小模型误报干扰，聚焦真实高风险行为",
            "复杂作业语义可理解并自动分级，支撑准实时预警处置",
        ],
        charts["ui_risk"],
    )

    # 22 Closing
    s = prs.slides.add_slide(blank)
    page_header(s, "共性价值与下一步", "提效 · 提质 · 可控 · 可复制", section="收束")
    left_items = [
        ("提效", "问规程、查台账、写分析初稿、录缺陷——开口驱动系统"),
        ("提质", "回答有依据，报告结构统一，减少随意表述"),
        ("可控", "资料不足会说明；关键环节可确认；查数只读"),
        ("可复制", "换知识、接数据后，可推广到同类行业/客户场景"),
    ]
    for i, (t, b) in enumerate(left_items):
        top = Inches(1.5 + i * 1.2)
        rounded_rect(s, Inches(0.65), top, Inches(7.1), Inches(1.05), fill=C["white"], line=C["line"])
        accent_bar(s, Inches(0.65), top, width=Inches(0.1), height=Inches(1.05), color=C["teal"] if i % 2 == 0 else C["navy"])
        add_textbox(s, Inches(1.0), top + Inches(0.15), Inches(1.2), Inches(0.35), t, size=14, bold=True, color=C["navy"])
        add_textbox(s, Inches(2.3), top + Inches(0.2), Inches(5.2), Inches(0.65), b, size=12, color=C["muted"])
    rounded_rect(s, Inches(8.0), Inches(1.5), Inches(4.7), Inches(4.85), fill=C["white"], line=C["line"])
    add_textbox(s, Inches(8.3), Inches(1.75), Inches(4.2), Inches(0.4), "建议下一步", size=16, bold=True, color=C["navy"])
    nexts = [
        "选 1～2 个高频场景先用起来",
        "持续补齐本厂规程与案例",
        "对接真实台账（只读）",
        "用真实截图替换效果示意",
    ]
    tb = s.shapes.add_textbox(Inches(8.3), Inches(2.4), Inches(4.2), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for i, line in enumerate(nexts, 1):
        txt = f"{i}.  {line}"
        if first:
            set_run(tf.paragraphs[0].add_run(), txt, size=13, color=C["ink"])
            first = False
        else:
            add_para(tf, txt, size=13, color=C["ink"], space_before=14)
    footer(s, 22, total)

    # Prefer -new.pptx (user working copy); fall back to primary if locked
    alt = OUT_DIR / "AI应用解决方案成果积累和应用场景示例.pptx"
    for target in (alt, OUT_PPTX):
        try:
            prs.save(str(target))
            return target
        except PermissionError:
            continue
    raise PermissionError(f"Cannot save PPTX: both {alt.name} and {OUT_PPTX.name} are locked")


def main():
    charts = build_charts()
    out = build_ppt(charts)
    print(f"OK: {out}")


if __name__ == "__main__":
    main()
