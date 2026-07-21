# -*- coding: utf-8 -*-
"""Leadership-friendly PPT: problems & solutions (plain language, less jargon)."""
from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

C = {
    "bg": "F4F6F8",
    "white": "FFFFFF",
    "ink": "1A2332",
    "muted": "5B6B7C",
    "line": "D5DCE4",
    "navy": "1B3A4B",
    "teal": "0F766E",
    "amber": "B45309",
}

OUT_DIR = Path(__file__).resolve().parent
ASSETS = OUT_DIR / "_ppt_assets_leader"
ASSETS.mkdir(exist_ok=True)
OUT_PPTX = OUT_DIR / "AI应用开发典型成果与场景汇报-领导通俗版.pptx"

plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False


def hex_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_run(run, text, size=14, bold=False, color=C["ink"], font="Microsoft YaHei"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = hex_rgb(color)
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = parse_xml(
        f'<a:ea xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{font}"/>'
    )
    rPr.append(ea)


def add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=C["ink"], align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size=size, bold=bold, color=color)
    return box


def add_para(tf, text, size=13, bold=False, color=C["ink"], space_before=6, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    set_run(p.add_run(), text, size=size, bold=bold, color=color)
    return p


def fill_shape(shape, color_hex: str):
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(color_hex)
    shape.line.fill.background()


def rounded_rect(slide, left, top, width, height, fill=C["white"], line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_shape(shape, fill)
    if line:
        shape.line.color.rgb = hex_rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def set_slide_bg(slide, color_hex=C["bg"]):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    fill_shape(bg, color_hex)
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.4), color=C["teal"]):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_shape(s, color)
    return s


def page_header(slide, title: str, subtitle: str | None = None, section: str | None = None):
    set_slide_bg(slide)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
    fill_shape(strip, C["navy"])
    if section:
        add_textbox(slide, Inches(0.7), Inches(0.28), Inches(11.5), Inches(0.3), section, size=11, color=C["teal"], bold=True)
    accent_bar(slide, Inches(0.7), Inches(0.58 if section else 0.45), height=Inches(0.36))
    add_textbox(
        slide,
        Inches(0.95),
        Inches(0.5 if section else 0.38),
        Inches(11.5),
        Inches(0.45),
        title,
        size=26,
        bold=True,
        color=C["navy"],
    )
    if subtitle:
        add_textbox(slide, Inches(0.95), Inches(1.0), Inches(11.5), Inches(0.35), subtitle, size=13, color=C["muted"])


def footer(slide, page: int, total: int, light: bool = True):
    ink = C["muted"] if light else "94A3B8"
    add_textbox(slide, Inches(0.7), Inches(7.1), Inches(9), Inches(0.3), "AI 应用成果汇报  ·  领导通俗版（问题与解决方案）", size=10, color=ink)
    add_textbox(slide, Inches(11.2), Inches(7.1), Inches(1.5), Inches(0.3), f"{page} / {total}", size=10, color=ink, align=PP_ALIGN.RIGHT)


def chart_pain_to_value(path: Path):
    fig, ax = plt.subplots(figsize=(12.2, 4.6), dpi=160)
    ax.set_xlim(0, 12.2)
    ax.set_ylim(0, 4.6)
    ax.axis("off")
    fig.patch.set_facecolor("#F4F6F8")
    ax.set_facecolor("#F4F6F8")

    left = FancyBboxPatch((0.4, 0.7), 5.0, 3.3, boxstyle="round,pad=0.03,rounding_size=0.15", facecolor="#1B3A4B", linewidth=0)
    right = FancyBboxPatch((6.8, 0.7), 5.0, 3.3, boxstyle="round,pad=0.03,rounding_size=0.15", facecolor="#0F766E", linewidth=0)
    mid = FancyBboxPatch((5.2, 1.85), 1.8, 1.0, boxstyle="round,pad=0.02,rounding_size=0.2", facecolor="#B45309", linewidth=0)
    ax.add_patch(left)
    ax.add_patch(right)
    ax.add_patch(mid)
    ax.text(2.9, 3.5, "现场常见问题", ha="center", color="white", fontsize=15, fontweight="bold")
    for i, t in enumerate(["规程翻不全、案例找不着", "查台账要找人写查询", "分析靠人工拼材料", "报告录入慢、格式乱"]):
        ax.text(0.75, 2.85 - i * 0.45, "·  " + t, color="#E2E8F0", fontsize=12)
    ax.text(9.3, 3.5, "我们提供的解决办法", ha="center", color="white", fontsize=15, fontweight="bold")
    for i, t in enumerate(["开口就能问规程与案例", "用普通话直接查数据", "一键生成分析初稿", "报告自动整理成台账"]):
        ax.text(7.15, 2.85 - i * 0.45, "·  " + t, color="#CCFBF1", fontsize=12)
    ax.text(6.1, 2.35, "对应\n解决", ha="center", va="center", color="white", fontsize=12, fontweight="bold")
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def chart_scene_plain(path: Path):
    fig, ax = plt.subplots(figsize=(12.2, 4.8), dpi=160)
    ax.set_xlim(0, 12.2)
    ax.set_ylim(0, 4.8)
    ax.axis("off")
    fig.patch.set_facecolor("#F4F6F8")
    ax.set_facecolor("#F4F6F8")
    groups = [
        (0.4, 2.5, 3.7, 1.9, "#0F766E", "日常问询", ["AI问答：规程与台账一口问", "直接用说话查数据", "知识资料可统一管理"]),
        (4.25, 2.5, 3.7, 1.9, "#1B3A4B", "分析研判", ["超温 / 检修策略 / 四管解读", "泄爆原因分方向梳理", "长报告可分段生成"]),
        (8.1, 2.5, 3.7, 1.9, "#0F766E", "现场与文档", ["拍照辅助缺陷研判", "检修报告自动整理入库", "事故描述可快速起草稿"]),
        (2.0, 0.3, 8.2, 1.8, "#334155", "希望带来的变化", ["少翻册、少找人、少手写查询", "报告更规范、结论更有依据", "关键环节可人工确认后再继续"]),
    ]
    for x, y, w, h, color, title, lines in groups:
        ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.02,rounding_size=0.12", facecolor=color, linewidth=0))
        ax.text(x + 0.25, y + h - 0.35, title, color="white", fontsize=13, fontweight="bold")
        for i, line in enumerate(lines):
            ax.text(x + 0.25, y + h - 0.75 - i * 0.38, "·  " + line, color="#E2E8F0", fontsize=11)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def chart_compare_plain(path: Path):
    fig, ax = plt.subplots(figsize=(11.5, 4.5), dpi=160)
    ax.set_xlim(0, 11.5)
    ax.set_ylim(0, 4.5)
    ax.axis("off")
    fig.patch.set_facecolor("#F4F6F8")
    ax.set_facecolor("#F4F6F8")
    ax.add_patch(FancyBboxPatch((0.4, 0.5), 5.0, 3.5, boxstyle="round,pad=0.03,rounding_size=0.12", facecolor="#64748B", linewidth=0))
    ax.add_patch(FancyBboxPatch((6.1, 0.5), 5.0, 3.5, boxstyle="round,pad=0.03,rounding_size=0.12", facecolor="#0F766E", linewidth=0))
    ax.text(2.9, 3.55, "只会「聊天」的 AI", ha="center", color="white", fontsize=14, fontweight="bold")
    ax.text(8.6, 3.55, "本平台（业务可用）", ha="center", color="white", fontsize=14, fontweight="bold")
    left_lines = ["容易凭感觉回答", "查不了厂内台账", "多轮对话容易跑偏", "难形成规范报告", "关键结论缺少把关"]
    right_lines = ["先找厂内资料再回答", "可用普通话查数据", "能听懂「它/第几点」", "可输出结构化分析稿", "拿不准时请人确认"]
    for i, (a, b) in enumerate(zip(left_lines, right_lines)):
        ax.text(0.75, 2.9 - i * 0.45, "·  " + a, color="#F1F5F9", fontsize=12)
        ax.text(6.45, 2.9 - i * 0.45, "·  " + b, color="#CCFBF1", fontsize=12)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def build_charts():
    paths = {
        "pain": ASSETS / "pain.png",
        "scenes": ASSETS / "scenes.png",
        "compare": ASSETS / "compare.png",
    }
    chart_pain_to_value(paths["pain"])
    chart_scene_plain(paths["scenes"])
    chart_compare_plain(paths["compare"])
    return paths


def build_ppt(charts: dict):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 12

    # Problem → Solution pairs (plain language, aligned to former 10 capabilities)
    pairs = [
        ("对话里说「它」「上次那个」，系统经常听不懂", "让系统能跟上对话上下文，听懂指代再回答"),
        ("想继续问「上面第 3 点」，又得把问题重说一遍", "自动记住回答要点，支持按条目接着问"),
        ("长规程切碎后，答案经常漏掉关键段落", "切分资料时保留前后衔接，降低漏答"),
        ("问到文字内容，相关图纸/附图对不上", "文字与图片一起管理，问到内容可带出附图"),
        ("查台账要找人写查询，错了还要反复改", "用普通话查数；写错了可自动修正，且只读不改库"),
        ("资料不够时，系统仍硬给出一个答案", "资料不足就再找、或先问清楚，不硬答"),
        ("做分析要到处拉数据，数据不全还硬写结论", "多路数据一起取；不够就说明缺口，不瞎编"),
        ("现场拍了照，还要回办公室查历史和规程", "看图 + 查历史 + 查规程一起做；拿不准可人工确认"),
        ("检修报告靠人工录入，慢且格式不统一", "长报告自动拆开整理，变成可入库的标准记录"),
        ("每换一个业务，又要重新搭一套说法和模板", "问答口径与报告模板可配置管理，便于推广复用"),
    ]

    # 1 Cover
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, C["navy"])
    deco = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), Inches(7.5))
    fill_shape(deco, C["teal"])
    add_textbox(s, Inches(1.1), Inches(2.0), Inches(11), Inches(0.4), "面向业务领导 · 通俗版", size=16, color="99F6E4", bold=True)
    add_textbox(s, Inches(1.1), Inches(2.55), Inches(11), Inches(1.0), "AI 应用：问题与解决方案", size=36, bold=True, color=C["white"])
    add_textbox(s, Inches(1.1), Inches(3.7), Inches(10.5), Inches(0.7), "我们解决了什么现场问题，又带来了哪些可用能力", size=18, color="CBD5E1")
    add_textbox(s, Inches(1.1), Inches(5.6), Inches(10), Inches(0.4), "能源电力工业场景  |  少用术语、多讲业务价值", size=13, color="94A3B8")
    footer(s, 1, total, light=False)

    # 2 Agenda
    s = prs.slides.add_slide(blank)
    page_header(s, "今天讲什么", "三条线看懂成果")
    agenda = [
        ("01", "现场有哪些痛点", "先对齐问题"),
        ("02", "我们怎么解决", "能力对应办法"),
        ("03", "能用在哪些场景", "效果怎么体现"),
        ("04", "和其他 AI 有何不同", "为何值得用"),
    ]
    for i, (num, title, desc) in enumerate(agenda):
        left = Inches(0.85 + i * 3.05)
        rounded_rect(s, left, Inches(2.0), Inches(2.85), Inches(3.5), fill=C["white"], line=C["line"])
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.95), Inches(2.4), Inches(0.9), Inches(0.9))
        fill_shape(circle, C["teal"] if i % 2 == 0 else C["navy"])
        add_textbox(s, left + Inches(0.95), Inches(2.58), Inches(0.9), Inches(0.55), num, size=18, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.15), Inches(3.55), Inches(2.55), Inches(0.5), title, size=15, bold=True, color=C["navy"], align=PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.2), Inches(4.2), Inches(2.45), Inches(0.8), desc, size=12, color=C["muted"], align=PP_ALIGN.CENTER)
    footer(s, 2, total)

    # 3 Pain overview
    s = prs.slides.add_slide(blank)
    page_header(s, "现场常见问题", "先看业务侧真正卡住的地方", section="问题")
    add_picture = lambda slide, path, left, top, width: slide.shapes.add_picture(str(path), left, top, width=width)
    add_picture(s, charts["pain"], Inches(0.55), Inches(1.5), Inches(12.2))
    footer(s, 3, total)

    # 4 What we deliver
    s = prs.slides.add_slide(blank)
    page_header(s, "我们交付了什么", "一句话：把「能问、能查、能分析、能整理」做成可用系统", section="解决方案总览")
    cards = [
        ("能问", "规程、案例、操作要点\n开口就能问"),
        ("能查", "台账、测点、缺陷记录\n用说话直接查"),
        ("能分析", "超温、泄爆、检修策略等\n自动形成分析初稿"),
        ("能整理", "检修报告自动整理\n变成标准记录"),
    ]
    for i, (t, b) in enumerate(cards):
        left = Inches(0.7 + i * 3.15)
        rounded_rect(s, left, Inches(2.0), Inches(3.0), Inches(3.8), fill=C["white"], line=C["line"])
        top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.0), Inches(3.0), Inches(1.1))
        fill_shape(top, C["teal"] if i % 2 == 0 else C["navy"])
        add_textbox(s, left, Inches(2.3), Inches(3.0), Inches(0.55), t, size=22, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.25), Inches(3.5), Inches(2.5), Inches(1.8), b, size=14, color=C["ink"], align=PP_ALIGN.CENTER)
    footer(s, 4, total)

    # 5 Four values plain
    s = prs.slides.add_slide(blank)
    page_header(s, "解决思路：四条原则", "不堆功能，先保证业务用得上、用得稳", section="解决方案总览")
    vals = [
        ("统一入口", "问答、查数、分析共用一套系统，避免每个项目另起炉灶"),
        ("有依据再答", "尽量先找厂内资料和数据，减少「凭感觉回答」"),
        ("拿不准可确认", "关键信息不确定时，可先请人确认再继续"),
        ("可复制推广", "换厂、换专题，主要通过配置调整，而不是重做一遍"),
    ]
    for i, (t, b) in enumerate(vals):
        top = Inches(1.55 + i * 1.2)
        rounded_rect(s, Inches(0.7), top, Inches(11.9), Inches(1.05), fill=C["white"], line=C["line"])
        accent_bar(s, Inches(0.7), top, width=Inches(0.12), height=Inches(1.05), color=C["teal"] if i % 2 == 0 else C["navy"])
        add_textbox(s, Inches(1.1), top + Inches(0.18), Inches(3.2), Inches(0.4), t, size=16, bold=True, color=C["navy"])
        add_textbox(s, Inches(4.5), top + Inches(0.25), Inches(7.7), Inches(0.55), b, size=14, color=C["muted"])
    footer(s, 5, total)

    # 6 Solutions 1-5
    s = prs.slides.add_slide(blank)
    page_header(s, "典型问题与解决办法（1/2）", "对应平台已沉淀的核心能力", section="解决方案明细")
    for i, (prob, sol) in enumerate(pairs[:5]):
        top = Inches(1.4 + i * 1.05)
        rounded_rect(s, Inches(0.55), top, Inches(6.0), Inches(0.95), fill=C["white"], line=C["line"])
        rounded_rect(s, Inches(6.75), top, Inches(6.0), Inches(0.95), fill=C["white"], line=C["line"])
        badge_l = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), top + Inches(0.28), Inches(0.7), Inches(0.4))
        fill_shape(badge_l, C["amber"])
        add_textbox(s, Inches(0.7), top + Inches(0.32), Inches(0.7), Inches(0.32), "问题", size=11, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(1.55), top + Inches(0.28), Inches(4.8), Inches(0.5), prob, size=12, color=C["ink"])
        badge_r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), top + Inches(0.28), Inches(0.7), Inches(0.4))
        fill_shape(badge_r, C["teal"])
        add_textbox(s, Inches(6.9), top + Inches(0.32), Inches(0.7), Inches(0.32), "办法", size=11, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(7.75), top + Inches(0.28), Inches(4.8), Inches(0.5), sol, size=12, color=C["ink"])
    footer(s, 6, total)

    # 7 Solutions 6-10
    s = prs.slides.add_slide(blank)
    page_header(s, "典型问题与解决办法（2/2）", "对应平台已沉淀的核心能力", section="解决方案明细")
    for i, (prob, sol) in enumerate(pairs[5:]):
        top = Inches(1.4 + i * 1.05)
        rounded_rect(s, Inches(0.55), top, Inches(6.0), Inches(0.95), fill=C["white"], line=C["line"])
        rounded_rect(s, Inches(6.75), top, Inches(6.0), Inches(0.95), fill=C["white"], line=C["line"])
        badge_l = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), top + Inches(0.28), Inches(0.7), Inches(0.4))
        fill_shape(badge_l, C["amber"])
        add_textbox(s, Inches(0.7), top + Inches(0.32), Inches(0.7), Inches(0.32), "问题", size=11, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(1.55), top + Inches(0.28), Inches(4.8), Inches(0.5), prob, size=12, color=C["ink"])
        badge_r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), top + Inches(0.28), Inches(0.7), Inches(0.4))
        fill_shape(badge_r, C["teal"])
        add_textbox(s, Inches(6.9), top + Inches(0.32), Inches(0.7), Inches(0.32), "办法", size=11, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(7.75), top + Inches(0.28), Inches(4.8), Inches(0.5), sol, size=12, color=C["ink"])
    footer(s, 7, total)

    # 8 Scenario map
    s = prs.slides.add_slide(blank)
    page_header(s, "已经能用在哪些地方", "按日常工作入口归集", section="场景与效果")
    add_picture(s, charts["scenes"], Inches(0.55), Inches(1.45), Inches(12.2))
    footer(s, 8, total)

    # 9 Scenario effects
    s = prs.slides.add_slide(blank)
    page_header(s, "核心场景与效果", "效果先用文字说明，后续可补截图", section="场景与效果")
    scenes = [
        ("AI问答", "规程问答 + 台账查数一口完成", "少翻册、少找人；统计类问题不必再等查询结果。"),
        ("超温 / 专项分析", "自动拉数并形成分析初稿", "缩短事故/异常初判时间；报告结构更统一。"),
        ("看图诊断", "现场拍照就能辅助研判", "边看图边对照历史与规程；拿不准可先确认再继续。"),
        ("检修报告整理", "Word/PDF 自动变成标准记录", "录入从「逐页抄」变成「上传后复核」，格式更统一。"),
        ("泄爆原因分析", "按方向梳理原因与防控建议", "事故调查起步更快，有规范草稿可讨论。"),
        ("知识资料管理", "规程案例统一入库可问可用", "资料不再散落；问答和分析都有依据可依。"),
    ]
    for i, (t, sub, eff) in enumerate(scenes):
        r, c = divmod(i, 3)
        left = Inches(0.65 + c * 4.15)
        top = Inches(1.5 + r * 2.6)
        rounded_rect(s, left, top, Inches(3.95), Inches(2.35), fill=C["white"], line=C["line"])
        accent_bar(s, left, top, width=Inches(0.1), height=Inches(2.35), color=C["teal"] if i % 2 == 0 else C["navy"])
        add_textbox(s, left + Inches(0.25), top + Inches(0.2), Inches(3.5), Inches(0.35), t, size=15, bold=True, color=C["navy"])
        add_textbox(s, left + Inches(0.25), top + Inches(0.6), Inches(3.5), Inches(0.4), sub, size=11, color=C["teal"])
        add_textbox(s, left + Inches(0.25), top + Inches(1.15), Inches(3.5), Inches(1.0), eff, size=12, color=C["muted"])
    footer(s, 9, total)

    # 10 Compare
    s = prs.slides.add_slide(blank)
    page_header(s, "和其他「聊天 AI」有何不同", "同样是 AI，但更贴近电厂日常工作", section="差异说明")
    add_picture(s, charts["compare"], Inches(0.9), Inches(1.55), Inches(11.5))
    footer(s, 10, total)

    # 11 Effect summary
    s = prs.slides.add_slide(blank)
    page_header(s, "希望带来的变化", "提效 · 提质 · 可控 · 可推广", section="收束")
    changes = [
        ("提效", "问规程、查台账、写分析初稿、录缺陷，从「人找系统」变成「开口驱动系统」"),
        ("提质", "回答尽量有资料依据，报告结构更规范，减少「拍脑袋」表述"),
        ("可控", "资料不足会说明；关键环节可人工确认；查询只读，不改业务库"),
        ("可推广", "同一套能力换知识、接数据后，可复制到同类电厂或场景"),
    ]
    for i, (t, b) in enumerate(changes):
        left = Inches(0.7 + (i % 2) * 6.2)
        top = Inches(1.6 + (i // 2) * 2.4)
        rounded_rect(s, left, top, Inches(5.9), Inches(2.15), fill=C["white"], line=C["line"])
        head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(5.9), Inches(0.7))
        fill_shape(head, C["teal"] if i % 2 == 0 else C["navy"])
        add_textbox(s, left, top + Inches(0.18), Inches(5.9), Inches(0.4), t, size=18, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.3), top + Inches(0.95), Inches(5.3), Inches(0.9), b, size=13, color=C["ink"])
    footer(s, 11, total)

    # 12 Next
    s = prs.slides.add_slide(blank)
    page_header(s, "建议下一步", "先抓高频场景，再逐步铺开", section="收束")
    nexts = [
        ("选 1～2 个高频场景先用起来", "例如 AI问答、超温分析或检修报告整理，先形成可感知效果"),
        ("把本厂规程与案例持续补齐", "资料越完整，回答越有依据，越少「答不上来」"),
        ("对接真实台账数据（只读）", "让「用说话查数」真正落到日常工作"),
        ("补齐演示截图与短视频", "便于对内汇报、对外交流时直观展示"),
    ]
    for i, (t, b) in enumerate(nexts):
        top = Inches(1.55 + i * 1.2)
        rounded_rect(s, Inches(0.7), top, Inches(11.9), Inches(1.05), fill=C["white"], line=C["line"])
        num = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), top + Inches(0.25), Inches(0.55), Inches(0.55))
        fill_shape(num, C["teal"] if i % 2 == 0 else C["navy"])
        add_textbox(s, Inches(1.0), top + Inches(0.35), Inches(0.55), Inches(0.35), str(i + 1), size=14, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(1.85), top + Inches(0.15), Inches(10.3), Inches(0.35), t, size=15, bold=True, color=C["navy"])
        add_textbox(s, Inches(1.85), top + Inches(0.55), Inches(10.3), Inches(0.4), b, size=13, color=C["muted"])
    footer(s, 12, total)

    try:
        prs.save(str(OUT_PPTX))
    except PermissionError:
        alt = OUT_DIR / "AI应用开发典型成果与场景汇报-领导通俗版-new.pptx"
        prs.save(str(alt))
        return alt
    return OUT_PPTX


def main():
    charts = build_charts()
    out = build_ppt(charts)
    print(f"OK: {out}")


if __name__ == "__main__":
    main()
