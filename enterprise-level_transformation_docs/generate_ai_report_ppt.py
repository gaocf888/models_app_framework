# -*- coding: utf-8 -*-
"""Generate conceptual company-report PPT for AI App Platform achievements."""
from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

# --- Visual system (industrial slate + teal; avoid purple / cream clichés) ---
C = {
    "bg": "F4F6F8",
    "white": "FFFFFF",
    "ink": "1A2332",
    "muted": "5B6B7C",
    "line": "D5DCE4",
    "navy": "1B3A4B",
    "teal": "0F766E",
    "teal_soft": "CCFBF1",
    "amber": "B45309",
    "amber_soft": "FEF3C7",
    "card": "FFFFFF",
    "accent2": "334155",
}

OUT_DIR = Path(__file__).resolve().parent
ASSETS = OUT_DIR / "_ppt_assets"
ASSETS.mkdir(exist_ok=True)
PPTX_PATH = OUT_DIR / "AI应用开发典型成果与场景汇报.pptx"
PPTX_PATH_ALT = OUT_DIR / "AI应用开发典型成果与场景汇报-v2.pptx"
PPTX_PATH_V3 = OUT_DIR / "AI应用开发典型成果与场景汇报-v3.pptx"

plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False


def hex_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_run(run, text, size=14, bold=False, color=C["ink"], font="Microsoft YaHei"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = hex_rgb(color)
    run.font.name = font
    # East Asian font
    rPr = run._r.get_or_add_rPr()
    ea = parse_xml(
        f'<a:ea xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{font}"/>'
    )
    rPr.append(ea)


def add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=C["ink"], align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, color=color)
    return box


def add_para(tf, text, size=13, bold=False, color=C["ink"], space_before=6, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    set_run(run, text, size=size, bold=bold, color=color)
    return p


def fill_shape(shape, color_hex: str):
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(color_hex)
    shape.line.fill.background()


def rounded_rect(slide, left, top, width, height, fill=C["white"], line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_shape(shape, fill)
    if line:
        shape.line.color.rgb = hex_rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # softer corners
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def set_slide_bg(slide, color_hex=C["bg"]):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    fill_shape(bg, color_hex)
    # send to back by moving XML element
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.4), color=C["teal"]):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_shape(s, color)
    return s


def page_header(slide, title: str, subtitle: str | None = None, section: str | None = None):
    set_slide_bg(slide)
    # top thin brand strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
    fill_shape(strip, C["navy"])
    if section:
        add_textbox(slide, Inches(0.7), Inches(0.28), Inches(11.5), Inches(0.3), section, size=11, color=C["teal"], bold=True)
    accent_bar(slide, Inches(0.7), Inches(0.58 if section else 0.45), height=Inches(0.36))
    add_textbox(
        slide,
        Inches(0.95),
        Inches(0.5 if section else 0.38),
        Inches(11.5),
        Inches(0.45),
        title,
        size=26,
        bold=True,
        color=C["navy"],
    )
    if subtitle:
        add_textbox(
            slide,
            Inches(0.95),
            Inches(1.0),
            Inches(11.5),
            Inches(0.35),
            subtitle,
            size=13,
            color=C["muted"],
        )


def footer(slide, page: int, total: int, light: bool = True):
    ink = C["muted"] if light else "94A3B8"
    add_textbox(
        slide,
        Inches(0.7),
        Inches(7.1),
        Inches(8),
        Inches(0.3),
        "大模型在线应用平台  ·  AI 应用成果汇报",
        size=10,
        color=ink,
    )
    add_textbox(
        slide,
        Inches(11.2),
        Inches(7.1),
        Inches(1.5),
        Inches(0.3),
        f"{page} / {total}",
        size=10,
        color=ink,
        align=PP_ALIGN.RIGHT,
    )


# ---------------- Charts ----------------
def chart_architecture(path: Path):
    fig, ax = plt.subplots(figsize=(12.5, 5.2), dpi=160)
    ax.set_xlim(0, 12.5)
    ax.set_ylim(0, 5.2)
    ax.axis("off")
    fig.patch.set_facecolor("#F4F6F8")
    ax.set_facecolor("#F4F6F8")

    layers = [
        (4.05, "#1B3A4B", "业务应用层", "AI问答 · 综合分析 · 看图诊断 · 检修抽取"),
        (3.05, "#0F766E", "编排层", "LangGraph · Prompt 治理 · HITL · 灰度回退"),
        (1.75, "#334155", "基座能力层", "混合 RAG · NL2SQL · LLM/多模态 · 多轮会话"),
        (0.55, "#64748B", "基础设施", "FastAPI · vLLM · EasySearch · Redis · MySQL · MinerU…"),
    ]
    for y, color, title, desc in layers:
        box = FancyBboxPatch(
            (0.8, y),
            10.9,
            0.85,
            boxstyle="round,pad=0.02,rounding_size=0.12",
            linewidth=0,
            facecolor=color,
        )
        ax.add_patch(box)
        ax.text(1.1, y + 0.52, title, color="white", fontsize=14, fontweight="bold", va="center")
        ax.text(1.1, y + 0.22, desc, color="#E2E8F0", fontsize=11, va="center")

    for y in (3.9, 2.9, 1.55):
        ax.annotate(
            "",
            xy=(6.25, y - 0.05),
            xytext=(6.25, y + 0.15),
            arrowprops=dict(arrowstyle="-|>", color="#94A3B8", lw=1.6),
        )
    ax.text(6.5, 3.95, "共享复用", color="#64748B", fontsize=9)
    ax.text(6.5, 2.95, "共享复用", color="#64748B", fontsize=9)
    ax.text(6.5, 1.6, "支撑", color="#64748B", fontsize=9)

    fig.tight_layout(pad=0.3)
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def chart_pillars(path: Path):
    fig, ax = plt.subplots(figsize=(12.2, 4.6), dpi=160)
    ax.set_xlim(0, 12.2)
    ax.set_ylim(0, 4.6)
    ax.axis("off")
    fig.patch.set_facecolor("#F4F6F8")
    ax.set_facecolor("#F4F6F8")

    items = [
        ("混合 RAG", "知识有据可依\n抑制幻觉"),
        ("NL2SQL", "自然语言问数\n安全只读"),
        ("LLM 推理", "流式 · 多模态\nPrompt 治理"),
        ("多轮会话", "热冷分层\n指代 · 大纲"),
        ("人机协同", "HITL 确认\n编排可回退"),
    ]
    colors = ["#0F766E", "#1B3A4B", "#0F766E", "#1B3A4B", "#0F766E"]
    w, gap = 2.05, 0.28
    x0 = 0.7
    for i, ((title, desc), color) in enumerate(zip(items, colors)):
        x = x0 + i * (w + gap)
        top = FancyBboxPatch(
            (x, 2.35), w, 1.55, boxstyle="round,pad=0.02,rounding_size=0.1", linewidth=0, facecolor=color
        )
        bottom = FancyBboxPatch(
            (x, 0.55),
            w,
            1.65,
            boxstyle="round,pad=0.02,rounding_size=0.1",
            linewidth=1.2,
            edgecolor="#D5DCE4",
            facecolor="white",
        )
        ax.add_patch(top)
        ax.add_patch(bottom)
        ax.text(x + w / 2, 3.15, title, ha="center", va="center", color="white", fontsize=13, fontweight="bold")
        ax.text(x + w / 2, 1.35, desc, ha="center", va="center", color="#1A2332", fontsize=11, linespacing=1.45)

    ax.text(6.1, 4.25, "可复用基座核心能力", ha="center", color="#1B3A4B", fontsize=15, fontweight="bold")
    fig.tight_layout(pad=0.2)
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def chart_dual_drive(path: Path):
    fig, ax = plt.subplots(figsize=(11.5, 4.2), dpi=160)
    ax.set_xlim(0, 11.5)
    ax.set_ylim(0, 4.2)
    ax.axis("off")
    fig.patch.set_facecolor("#F4F6F8")
    ax.set_facecolor("#F4F6F8")

    left = FancyBboxPatch((0.5, 0.9), 4.4, 2.6, boxstyle="round,pad=0.03,rounding_size=0.15", facecolor="#0F766E", linewidth=0)
    right = FancyBboxPatch((6.6, 0.9), 4.4, 2.6, boxstyle="round,pad=0.03,rounding_size=0.15", facecolor="#1B3A4B", linewidth=0)
    mid = FancyBboxPatch((4.55, 1.65), 2.4, 1.1, boxstyle="round,pad=0.02,rounding_size=0.2", facecolor="#B45309", linewidth=0)
    ax.add_patch(left)
    ax.add_patch(right)
    ax.add_patch(mid)
    ax.text(2.7, 2.85, "非结构化知识", ha="center", color="white", fontsize=15, fontweight="bold")
    ax.text(2.7, 2.15, "混合 RAG / GraphRAG\n规程 · 案例 · 图文", ha="center", color="#CCFBF1", fontsize=12)
    ax.text(8.8, 2.85, "结构化数据", ha="center", color="white", fontsize=15, fontweight="bold")
    ax.text(8.8, 2.15, "NL2SQL 安全问数\n台账 · 测点 · 缺陷", ha="center", color="#E2E8F0", fontsize=12)
    ax.text(5.75, 2.2, "双轮驱动", ha="center", va="center", color="white", fontsize=13, fontweight="bold")
    ax.text(5.75, 3.85, "相对「只包一层 Chat」的能力底座", ha="center", color="#5B6B7C", fontsize=11)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def chart_capability_radar(path: Path):
    import numpy as np

    labels = ["知识检索", "安全问数", "多轮理解", "工业分析", "看图诊断", "文档智能"]
    # Conceptual maturity scores for narrative (not literal metrics)
    ours = [92, 90, 88, 90, 87, 86]
    baseline = [45, 25, 40, 35, 40, 40]
    N = len(labels)
    angles = np.linspace(0, 2 * np.pi, N, endpoint=False).tolist()
    ours += ours[:1]
    baseline += baseline[:1]
    angles += angles[:1]

    fig, ax = plt.subplots(figsize=(6.8, 6.2), dpi=160, subplot_kw=dict(polar=True))
    fig.patch.set_facecolor("#F4F6F8")
    ax.set_facecolor("#F4F6F8")
    ax.plot(angles, baseline, color="#94A3B8", linewidth=1.5, label="单点 LLM 封装")
    ax.fill(angles, baseline, color="#94A3B8", alpha=0.15)
    ax.plot(angles, ours, color="#0F766E", linewidth=2.2, label="本平台")
    ax.fill(angles, ours, color="#0F766E", alpha=0.22)
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(labels, fontsize=11, color="#1A2332")
    ax.set_yticks([20, 40, 60, 80])
    ax.set_yticklabels(["", "", "", ""], color="#94A3B8")
    ax.set_ylim(0, 100)
    ax.spines["polar"].set_color("#D5DCE4")
    ax.grid(color="#CBD5E1", linestyle="--", linewidth=0.7)
    ax.legend(loc="upper right", bbox_to_anchor=(1.28, 1.12), fontsize=10, frameon=False)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def chart_scenario_map(path: Path):
    fig, ax = plt.subplots(figsize=(12.2, 5.0), dpi=160)
    ax.set_xlim(0, 12.2)
    ax.set_ylim(0, 5.0)
    ax.axis("off")
    fig.patch.set_facecolor("#F4F6F8")
    ax.set_facecolor("#F4F6F8")

    groups = [
        (0.4, 2.7, 3.7, 1.9, "#0F766E", "交互入口", ["AI问答", "独立 NL2SQL 问数", "知识库运维"]),
        (4.25, 2.7, 3.7, 1.9, "#1B3A4B", "分析决策", ["超温 / 检修策略 / 四管", "泄爆分析", "综合分析智能体"]),
        (8.1, 2.7, 3.7, 1.9, "#0F766E", "现场与文档", ["看图诊断（缺陷/泄爆）", "检修报告结构化提取", "结构化报告输出"]),
        (2.2, 0.35, 7.8, 1.85, "#334155", "共性效果", ["提效：自然语言驱动系统", "提质：有据可依 · 结构可控", "可控：可确认 · 可续问 · 可扩展"]),
    ]
    for x, y, w, h, color, title, lines in groups:
        box = FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.02,rounding_size=0.12", facecolor=color, linewidth=0)
        ax.add_patch(box)
        ax.text(x + 0.25, y + h - 0.35, title, color="white", fontsize=13, fontweight="bold")
        for i, line in enumerate(lines):
            ax.text(x + 0.25, y + h - 0.75 - i * 0.38, "·  " + line, color="#E2E8F0", fontsize=11)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def chart_effect_bars(path: Path):
    fig, ax = plt.subplots(figsize=(11.2, 4.4), dpi=160)
    fig.patch.set_facecolor("#F4F6F8")
    labels = ["信息获取\n效率", "报告初稿\n产出", "缺陷录入\n规范性", "结论\n可追溯", "现场处置\n辅助"]
    values = [85, 80, 88, 82, 78]
    y = range(len(labels))
    bars = ax.barh(list(y), values, color=["#0F766E", "#1B3A4B", "#0F766E", "#1B3A4B", "#0F766E"], height=0.55)
    ax.set_yticks(list(y))
    ax.set_yticklabels(labels, fontsize=12, color="#1A2332")
    ax.set_xlim(0, 100)
    ax.set_xlabel("概念示意 · 业务价值强度（非实测指标）", fontsize=10, color="#5B6B7C")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#D5DCE4")
    ax.spines["bottom"].set_color("#D5DCE4")
    ax.tick_params(colors="#5B6B7C")
    ax.set_facecolor("#F4F6F8")
    for bar, v in zip(bars, values):
        ax.text(v + 1.5, bar.get_y() + bar.get_height() / 2, f"{v}", va="center", color="#1B3A4B", fontsize=11, fontweight="bold")
    ax.set_title("场景效果维度示意", fontsize=14, color="#1B3A4B", fontweight="bold", pad=12)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


def build_charts():
    paths = {
        "arch": ASSETS / "arch.png",
        "pillars": ASSETS / "pillars.png",
        "dual": ASSETS / "dual.png",
        "radar": ASSETS / "radar.png",
        "scenes": ASSETS / "scenes.png",
        "effects": ASSETS / "effects.png",
    }
    chart_architecture(paths["arch"])
    chart_pillars(paths["pillars"])
    chart_dual_drive(paths["dual"])
    chart_capability_radar(paths["radar"])
    chart_scenario_map(paths["scenes"])
    chart_effect_bars(paths["effects"])
    return paths


def add_picture(slide, path: Path, left, top, width):
    slide.shapes.add_picture(str(path), left, top, width=width)


def card(slide, left, top, width, height, title, body, accent=C["teal"]):
    rounded_rect(slide, left, top, width, height, fill=C["white"], line=C["line"])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    fill_shape(bar, accent)
    tb = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.15), width - Inches(0.35), height - Inches(0.25))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_run(p.add_run(), title, size=14, bold=True, color=C["navy"])
    add_para(tf, body, size=12, color=C["muted"], space_before=8)


def build_ppt(charts: dict):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 14

    # Shared capability catalog (overview tagline + brief explanation)
    caps = [
        (
            "指代消解 P0～P3",
            "「它 / 上次那个」可绑定续问",
            "规则判型后融合检索 query，并可叠加对话锚块、会话槽位与灰区 Coref，让多轮弱指代仍能答到正确实体与话题。",
        ),
        (
            "章节大纲续问",
            "「上面第 N 点」结构化回锚",
            "助手回答后异步抽取条目大纲并保存；用户按「第几点」追问时可精确回锚对应要点，与指代消解互补。",
        ),
        (
            "滑动窗口切块",
            "长规程连贯，降低切缝漏召",
            "知识摄入按结构/语义切分，超长段用 overlap 滑动窗口保留相邻上下文，缓解答案落在切缝导致的漏召。",
        ),
        (
            "图文关联召回",
            "问文字可带出附图依据",
            "正文与图块双向关联入库；正文命中后可扩展召回附图，支撑规程/案例类「问文字、见附图」。",
        ),
        (
            "NL2SQL 双闭环",
            "校验/执行失败可 refine",
            "生成校验失败与 EXPLAIN/执行失败均可自动修正 SQL，且修正后仍须过完整安全校验，不能绕过只读边界。",
        ),
        (
            "C-RAG 质量闭环",
            "低质召回不硬答，可重写澄清",
            "对检索质量做门控：不足则重写 query、有限次重试；仍不达标则转澄清，避免低质量上下文硬答。",
        ),
        (
            "分析并行 + 质量门",
            "多维取数，不足则降级",
            "分析计划按依赖分层并行 NL2SQL 取数；数据不足时降级提示而非编造，保障报告可用性与可信度。",
        ),
        (
            "看图三臂 + HITL",
            "视觉 ‖ 问数 ‖ 知识可确认",
            "视觉判定、历史问数、业务知识三路并行后合成；机组/范围不确定时可人机确认，降低取错数风险。",
        ),
        (
            "检修分块并行抽取",
            "三阶段 LLM，长文档结构化",
            "长检修报告分块后走 parse→classify→repair 三阶段并行 LLM 抽取，输出可入库的标准缺陷结构化数据。",
        ),
        (
            "Prompt 配置治理",
            "版本可控 · 场景可扩展",
            "Prompt 与报告规格集中配置、可按场景/版本分流；新专题优先加配置扩展，少改核心引擎。",
        ),
    ]

    # 1 Cover
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, C["navy"])
    # decorative teal block
    deco = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), Inches(7.5))
    fill_shape(deco, C["teal"])
    add_textbox(s, Inches(1.1), Inches(2.0), Inches(11), Inches(0.4), "企业级 AI 应用成果汇报", size=16, color="99F6E4", bold=True)
    add_textbox(s, Inches(1.1), Inches(2.55), Inches(11), Inches(1.0), "大模型在线应用平台", size=40, bold=True, color=C["white"])
    add_textbox(
        s,
        Inches(1.1),
        Inches(3.7),
        Inches(10.5),
        Inches(0.8),
        "成果积累  ·  使用场景与效果分享",
        size=20,
        color="CBD5E1",
    )
    add_textbox(
        s,
        Inches(1.1),
        Inches(5.6),
        Inches(10),
        Inches(0.4),
        "面向能源电力工业场景  |  基座能力 + 业务应用",
        size=13,
        color="94A3B8",
    )
    footer(s, 1, total, light=False)

    # 2 Agenda
    s = prs.slides.add_slide(blank)
    page_header(s, "汇报结构", "两条主线，概念先行")
    agenda = [
        ("01", "定位与价值", "平台是什么，解决什么问题"),
        ("02", "成果积累", "基座能力 · 关键技术 · 简要说明"),
        ("03", "场景与效果", "已落地应用与文字效果描述"),
        ("04", "收束展望", "可复制、可演示、可深化"),
    ]
    for i, (num, title, desc) in enumerate(agenda):
        left = Inches(0.85 + i * 3.05)
        rounded_rect(s, left, Inches(2.0), Inches(2.85), Inches(3.6), fill=C["white"], line=C["line"])
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.95), Inches(2.35), Inches(0.9), Inches(0.9))
        fill_shape(circle, C["teal"] if i % 2 == 0 else C["navy"])
        add_textbox(s, left + Inches(0.95), Inches(2.52), Inches(0.9), Inches(0.55), num, size=18, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.15), Inches(3.5), Inches(2.55), Inches(0.45), title, size=16, bold=True, color=C["navy"], align=PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.2), Inches(4.15), Inches(2.45), Inches(1.0), desc, size=12, color=C["muted"], align=PP_ALIGN.CENTER)
    footer(s, 2, total)

    # 3 Positioning
    s = prs.slides.add_slide(blank)
    page_header(s, "平台定位", "一句话说清我们在做什么")
    rounded_rect(s, Inches(0.7), Inches(1.55), Inches(11.9), Inches(1.35), fill=C["white"], line=C["line"])
    add_textbox(
        s,
        Inches(1.1),
        Inches(1.8),
        Inches(11.3),
        Inches(0.9),
        "面向能源电力工业场景的企业级 AI 应用基座：统一沉淀推理、检索、问数与多轮会话能力，并交付开箱即用的业务应用。",
        size=16,
        color=C["ink"],
    )
    values = [
        ("统一基座", "推理 · 检索 · 问数 · 会话\n收敛为一套服务"),
        ("生产可控", "只读 SQL · 知识域隔离\nPrompt 治理 · HITL"),
        ("场景可扩展", "新业务优先配置接入\n少改底层引擎"),
        ("双轮驱动", "非结构化知识 + 结构化数据\n而非「只包一层 Chat」"),
    ]
    for i, (t, b) in enumerate(values):
        left = Inches(0.7 + (i % 4) * 3.1)
        card(s, left, Inches(3.25), Inches(2.95), Inches(2.7), t, b, accent=C["teal"] if i % 2 == 0 else C["navy"])
    footer(s, 3, total)

    # 4 Architecture
    s = prs.slides.add_slide(blank)
    page_header(s, "成果总览：分层架构", "可复用基座 + 可交付应用 + 可复制部署", section="PART 01  ·  成果积累")
    add_picture(s, charts["arch"], Inches(0.55), Inches(1.55), Inches(12.2))
    footer(s, 4, total)

    # 5 Dual drive + pillars
    s = prs.slides.add_slide(blank)
    page_header(s, "基座能力：双轮驱动", "知识检索与结构化问数并列", section="PART 01  ·  成果积累")
    add_picture(s, charts["dual"], Inches(0.9), Inches(1.55), Inches(11.5))
    add_textbox(
        s,
        Inches(0.9),
        Inches(6.35),
        Inches(11.5),
        Inches(0.4),
        "积累形态：不是零散 Demo，而是可复用、可交付、可复制的三位一体能力。",
        size=13,
        color=C["muted"],
        align=PP_ALIGN.CENTER,
    )
    footer(s, 5, total)

    # 6 Five pillars
    s = prs.slides.add_slide(blank)
    page_header(s, "基座核心能力", "上层场景共享复用的核心能力", section="PART 01  ·  成果积累")
    add_picture(s, charts["pillars"], Inches(0.55), Inches(1.55), Inches(12.2))
    footer(s, 6, total)

    # 7 Key capabilities overview
    s = prs.slides.add_slide(blank)
    page_header(s, "企业级关键技术能力（精选）", "相对通用 LLM 封装的差异化深化", section="PART 01  ·  成果积累")
    for i, (t, tag, _desc) in enumerate(caps):
        r, c = divmod(i, 5)
        left = Inches(0.65 + c * 2.5)
        top = Inches(1.55 + r * 2.5)
        rounded_rect(s, left, top, Inches(2.35), Inches(2.2), fill=C["white"], line=C["line"])
        num = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.85), top + Inches(0.25), Inches(0.55), Inches(0.55))
        fill_shape(num, C["teal"] if i % 2 == 0 else C["navy"])
        add_textbox(s, left + Inches(0.85), top + Inches(0.33), Inches(0.55), Inches(0.4), f"{i+1:02d}", size=12, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.12), top + Inches(0.95), Inches(2.1), Inches(0.55), t, size=12, bold=True, color=C["navy"], align=PP_ALIGN.CENTER)
        add_textbox(s, left + Inches(0.12), top + Inches(1.45), Inches(2.1), Inches(0.55), tag, size=11, color=C["muted"], align=PP_ALIGN.CENTER)
    footer(s, 7, total)

    # 8 Capability briefs 01-05
    s = prs.slides.add_slide(blank)
    page_header(s, "关键能力简要说明（1/2）", "精选能力 01～05：解决什么问题", section="PART 01  ·  成果积累")
    for i, (t, tag, desc) in enumerate(caps[:5]):
        top = Inches(1.45 + i * 1.05)
        rounded_rect(s, Inches(0.65), top, Inches(12.0), Inches(0.95), fill=C["white"], line=C["line"])
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), top + Inches(0.22), Inches(0.7), Inches(0.5))
        fill_shape(badge, C["teal"] if i % 2 == 0 else C["navy"])
        try:
            badge.adjustments[0] = 0.2
        except Exception:
            pass
        add_textbox(s, Inches(0.85), top + Inches(0.3), Inches(0.7), Inches(0.35), f"{i+1:02d}", size=12, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(1.75), top + Inches(0.12), Inches(4.2), Inches(0.35), t, size=14, bold=True, color=C["navy"])
        add_textbox(s, Inches(6.0), top + Inches(0.14), Inches(6.3), Inches(0.3), tag, size=11, color=C["teal"])
        add_textbox(s, Inches(1.75), top + Inches(0.48), Inches(10.5), Inches(0.4), desc, size=12, color=C["muted"])
    footer(s, 8, total)

    # 9 Capability briefs 06-10
    s = prs.slides.add_slide(blank)
    page_header(s, "关键能力简要说明（2/2）", "精选能力 06～10：解决什么问题", section="PART 01  ·  成果积累")
    for i, (t, tag, desc) in enumerate(caps[5:]):
        idx = i + 6
        top = Inches(1.45 + i * 1.05)
        rounded_rect(s, Inches(0.65), top, Inches(12.0), Inches(0.95), fill=C["white"], line=C["line"])
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), top + Inches(0.22), Inches(0.7), Inches(0.5))
        fill_shape(badge, C["teal"] if idx % 2 == 0 else C["navy"])
        try:
            badge.adjustments[0] = 0.2
        except Exception:
            pass
        add_textbox(s, Inches(0.85), top + Inches(0.3), Inches(0.7), Inches(0.35), f"{idx:02d}", size=12, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(1.75), top + Inches(0.12), Inches(4.2), Inches(0.35), t, size=14, bold=True, color=C["navy"])
        add_textbox(s, Inches(6.0), top + Inches(0.14), Inches(6.3), Inches(0.3), tag, size=11, color=C["teal"])
        add_textbox(s, Inches(1.75), top + Inches(0.48), Inches(10.5), Inches(0.4), desc, size=12, color=C["muted"])
    footer(s, 9, total)

    # 10 Differentiation radar
    s = prs.slides.add_slide(blank)
    page_header(s, "差异化：平台 vs 单点封装", "概念对比示意，突出能力广度与深度", section="PART 01  ·  成果积累")
    add_picture(s, charts["radar"], Inches(0.4), Inches(1.35), Inches(6.3))
    # right side key points
    points = [
        "知识问答：混合 RAG + 本厂隔离 + 指代",
        "数据查询：安全只读 NL2SQL，非手写 SQL",
        "多轮对话：热冷会话 + 大纲结构化续问",
        "工业分析：并行取数 + 结构化报告输出",
        "看图诊断：视觉 + 台账 + 规程三路合成",
        "文档智能：分块并行三阶段结构化抽取",
    ]
    rounded_rect(s, Inches(6.9), Inches(1.7), Inches(5.7), Inches(4.7), fill=C["white"], line=C["line"])
    add_textbox(s, Inches(7.2), Inches(1.95), Inches(5.2), Inches(0.4), "关键差异一览", size=16, bold=True, color=C["navy"])
    tb = s.shapes.add_textbox(Inches(7.2), Inches(2.5), Inches(5.2), Inches(3.6))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for ptxt in points:
        if first:
            set_run(tf.paragraphs[0].add_run(), "●  " + ptxt, size=13, color=C["ink"])
            first = False
        else:
            add_para(tf, "●  " + ptxt, size=13, color=C["ink"], space_before=12)
    footer(s, 10, total)

    # 11 Section divider scenarios
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, C["teal"])
    deco = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), Inches(7.5))
    fill_shape(deco, C["navy"])
    add_textbox(s, Inches(1.1), Inches(2.6), Inches(11), Inches(0.4), "PART 02", size=16, bold=True, color="99F6E4")
    add_textbox(s, Inches(1.1), Inches(3.15), Inches(11), Inches(0.8), "使用场景及效果分享", size=36, bold=True, color=C["white"])
    add_textbox(s, Inches(1.1), Inches(4.15), Inches(10), Inches(0.5), "效果以文字描述为主，后续可扩展应用截图 / 演示素材", size=15, color="CCFBF1")
    footer(s, 11, total, light=False)

    # 12 Scenario map
    s = prs.slides.add_slide(blank)
    page_header(s, "场景全景", "已落地应用按三类入口归集", section="PART 02  ·  场景与效果")
    add_picture(s, charts["scenes"], Inches(0.55), Inches(1.5), Inches(12.2))
    footer(s, 12, total)

    # 13 Scenario effects cards
    s = prs.slides.add_slide(blank)
    page_header(s, "核心场景与效果", "概念级效果描述，便于后续配图", section="PART 02  ·  场景与效果")
    scenes = [
        ("AI问答", "规程问答 + 台账查数一站式入口", "7×24 自然语言入口；少翻册、少写 SQL；支持「它 / 第几点」续问。"),
        ("超温 / 专项分析", "并行取数生成结构化分析报告", "缩短信息收集与初判时间；报告结构统一；数据不足可降级提示。"),
        ("看图诊断", "现场拍照 + 历史数据 + 规程知识", "巡检/停机「拍即诊」；风险与处置建议可讨论；范围可人确认。"),
        ("检修报告抽取", "Word/PDF → 标准缺陷 JSON", "录入由抄录转为复核；格式统一可入库；长文档分块并行抽取。"),
        ("泄爆溯源分析", "五类×三层结构化溯源", "事故调查有规范草稿起步；案例与规程可引用、可复查。"),
        ("知识库运维", "文档摄入 · 域隔离 · 检索调优", "知识资产可运营；支撑客服与分析「答得准、不串域」。"),
    ]
    for i, (t, sub, eff) in enumerate(scenes):
        r, c = divmod(i, 3)
        left = Inches(0.65 + c * 4.15)
        top = Inches(1.5 + r * 2.6)
        rounded_rect(s, left, top, Inches(3.95), Inches(2.35), fill=C["white"], line=C["line"])
        accent_bar(s, left, top, width=Inches(0.1), height=Inches(2.35), color=C["teal"] if i % 2 == 0 else C["navy"])
        add_textbox(s, left + Inches(0.25), top + Inches(0.2), Inches(3.5), Inches(0.35), t, size=15, bold=True, color=C["navy"])
        add_textbox(s, left + Inches(0.25), top + Inches(0.6), Inches(3.5), Inches(0.4), sub, size=11, color=C["teal"])
        add_textbox(s, left + Inches(0.25), top + Inches(1.15), Inches(3.5), Inches(1.0), eff, size=12, color=C["muted"])
    footer(s, 13, total)

    # 14 Closing with effect chart
    s = prs.slides.add_slide(blank)
    page_header(s, "效果共性与下一步", "提效 · 提质 · 可控 · 可复制", section="收束")
    add_picture(s, charts["effects"], Inches(0.45), Inches(1.45), Inches(7.0))
    rounded_rect(s, Inches(7.7), Inches(1.6), Inches(4.9), Inches(4.7), fill=C["white"], line=C["line"])
    add_textbox(s, Inches(8.0), Inches(1.85), Inches(4.4), Inches(0.4), "下一步建议", size=16, bold=True, color=C["navy"])
    nexts = [
        "选定 1～2 个高频场景做效果量化",
        "完善本厂知识域运营与更新机制",
        "按场景清单补截图 / 演示素材",
        "现场复制：接库 · 接知识 · 开场景",
    ]
    tb = s.shapes.add_textbox(Inches(8.0), Inches(2.5), Inches(4.4), Inches(3.4))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for i, line in enumerate(nexts, 1):
        txt = f"{i}.  {line}"
        if first:
            set_run(tf.paragraphs[0].add_run(), txt, size=13, color=C["ink"])
            first = False
        else:
            add_para(tf, txt, size=13, color=C["ink"], space_before=14)
    footer(s, 14, total)

    outputs = []
    for path in (PPTX_PATH, PPTX_PATH_ALT, PPTX_PATH_V3):
        try:
            prs.save(str(path))
            outputs.append(path)
        except PermissionError:
            print(f"SKIP (locked): {path}")
    if not outputs:
        raise PermissionError("Cannot write any PPTX; close open files and retry")
    return outputs[0]


def main():
    charts = build_charts()
    out = build_ppt(charts)
    print(f"OK: {out}")


if __name__ == "__main__":
    main()
